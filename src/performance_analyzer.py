from src import common_module as cm
from src.common.constants import SystemConstants, DbTypeConstants
from src.common.utils import DateUtils, SystemUtils, SqlUtils
from src.analysis_extend_target import OracleTarget
from src.ppt.ppt_writer import SlideManager

import pandas as pd
import re

from resources.config_manager import Config

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Pt
from pptx.enum.chart import XL_DATA_LABEL_POSITION
from pptx.enum.chart import XL_MARKER_STYLE
from pptx.util import Inches
from sql.common_sql import PptSql
from datetime import datetime


class PerformanceAnalyzer(cm.CommonModule):
    """
    ppt_report
    """

    FONT_NAME = "나눔스퀘어 네오 Light"

    SMALL_FONT_SIZE = Pt(8)
    DEFAULT_FONT_SIZE = Pt(10)
    MEDIUM_FONT_SIZE = Pt(11)
    LARGE_FONT_SIZE = Pt(24)

    BASE_LINE_SPACE = 1.0
    LONG_LINE_SPACE = 1.5

    BOLD_TRUE = True
    BOLD_FALSE = False

    YELLOW_COLOR = RGBColor(255, 217, 102)
    WHILE_COLOR = RGBColor(255, 255, 255)
    LIGHT_GRAY_COLOR = RGBColor(240, 240, 240)
    GREEN_COLOR = RGBColor(112, 173, 71)
    LIGHT_RED_COLOR = RGBColor(198, 29, 81)
    LIGHT_GREEN_COLOR = RGBColor(25, 184, 97)
    BLUE_GRAY_COLOR = RGBColor(65, 84, 110)
    DARK_GRAY_COLOR = RGBColor(89, 89, 89)
    LIGHT_BLUE_COLOR = RGBColor(61, 129, 246)
    BLACK_COLOR = RGBColor(0, 0, 0)

    ACTIVITY_INFO = "#ACTIVITY_INFO"
    SCHEDULE = "#SCHEDULE"
    LOAD_PROFILE = "#LOAD_PROFILE"
    TIME_MODEL = "#TIME_MODEL"
    DB_SYSTEM = "#DB_SYSTEM"
    TON_N = "#TOP_N_WAIT_EVENTS"
    TOP_SCHEMA_SQL = "#TOP_SCHEMA_SQL"
    LITERAL_SQL = "#LITERAL_SQL"
    METRIC = "#METRIC"
    RAC = "#RAC"
    MEMORY = "#MEMORY"
    TOP_3 = "#TOP_3_WAIT_EVENTS"
    TOP_SEG_IO = "#TOP_SEO_IO"
    AUTO_TASK = "#AUTO_TASK"
    OPTIMIZER_STATISTICS = "#OPTIMIZER_STATISTICS"
    OPTIMIZER_STATISTICS_ZERO_NULL = "#OPTIMIZER_STATISTICS_ZERO_NULL"
    MEMORY_SGA = "#MEMORY_ADVICE_SGA"
    MEMORY_PGA = "#MEMORY_ADVICE_PGA"
    MEMORY_SHARED_POOL = "#MEMORY_ADVICE_SHARED_POOL"
    MEMORY_DB_BUFFER_CACHE = "#MEMORY_ADVICE_DB_BUFFER_CACHE"
    DATABASE_PARAMETER = "#DATABASE_PARAMETER"
    TOP_LITERAL_RATIO = "#TOP_LITERAL_SQL_RATIO"
    TOP_SEGMENTS_BLOCK = "#TOP_SEGMENTS_BLOCK"
    TOP_SEGMENTS_RATIO = "#TOP_SEGMENTS_RATIO"

    COLOR_LIST = [
        RGBColor(65, 84, 110),
        RGBColor(198, 29, 81),
        RGBColor(112, 173, 71),
        RGBColor(255, 217, 102),
        RGBColor(216, 48, 196),
    ]

    y_axis_dict = {
        "Host CPU Utilization (%)": "Usage(%)",
        "Logical Reads Per Sec": "Block (K=1000)",
        "Physical Reads Per Sec": "Block (K=1000)",
    }

    metric_top_menu_dict = {
        "Host CPU Utilization (%)": "CPU Usage",
        "Average Active Sessions": "Average Active Sessions",
        "Executions Per Sec": "SQL Execution",
        "Logical Reads Per Sec": "Logical Reads",
        "Physical Reads Per Sec": "Physical Reads",
        "Hard Parse Count Per Sec": "Hard Parse",
    }

    def __init__(self, logger):
        super().__init__(logger=logger)
        self.ot: OracleTarget = None
        self.sql_path = None
        self.position = None
        self.presentiton_path = None

        self.instance_number = None
        self.instance_name = None

    def main_process(self):
        """main_process"""
        self.logger.info("performance_analyzer.pptx")
        self._insert_extend_target_data()

        self.presentiton_path = f"{self.config['home']}/{SystemConstants.TEMPLATE_PATH}/template.pptx"
        self.presentation = Presentation(self.presentiton_path)
        self.instance_number = self._extract_instance_info("instance_number")
        self.instance_name = self._extract_instance_info("instance_name")
        self.oracle_version = self._extract_instance_info("oracle_version")

        self.today_date = datetime.today().strftime("%Y. %m. %d")
        self.period = self._extract_period()

        self.sql_path = f"{self.config['home']}/" + SystemConstants.CHART_SQL
        self.position = Config("position").get_config()

        self._execute_first_slide()
        self._execute_activity_info()
        self._execute_schedule()
        self._execute_load_profile()
        self._execute_time_model()
        self._execute_db_system()
        self._execute_metric()
        self._execute_rac()
        self._execute_top_n_wait_events()
        self._execute_top_3_wait_events()
        self._execute_top_schema_sql()
        self._execute_memory_advice_sga()
        self._execute_memory_advice_pga()
        self._execute_memory_shared_pool()
        self._execute_memory_db_buffer_cache()
        self._execute_database_parameter()
        self.top_literal_sql_ratio()
        self._execute_literal_sql()
        self._execute_auto_task()
        self._execute_optimizer_statistics()
        self._execute_optimizer_statistics_zero_null()
        self._execute_top_seg_io_block()
        self._execute_top_seg_io_ratio()
        self._execute_top_seg_io_table()
        SlideManager.delete_slide(self.presentation)
        self.presentation.save(f"{self.config['home']}/{SystemConstants.EXPORT_DIR}/analyzer_report.pptx")

    def _execute_delete_slide(self):
        """
        #으로 시작하는 슬라이드 삭제
        """
        xml_slides = self.presentation.slides._sldIdLst
        xml_slides_list = list(xml_slides)
        for idx, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.has_text_frame and re.search(r"^#", shape.text_frame.text):
                    xml_slides.remove(xml_slides_list[idx])

    def _insert_extend_target_data(self):
        """
        DB 확장 타겟(실제 분석 대상 DB) 데이터 저장 함수
        """
        extend_target_repo_list = self.config["maxgauge_repo"].get("extend_target_repo", [])

        for extend_target_repo in extend_target_repo_list:
            extend_target_repo["analysis_target_type"] = self.config["maxgauge_repo"]["analysis_target_type"]  # oracle

            if str(self.config["maxgauge_repo"]["analysis_target_type"]).lower() == DbTypeConstants.ORACLE:
                if self.ot is None:
                    self.ot = OracleTarget(self.logger, self.config)

                self.ot.set_extend_target_config(extend_target_repo)
                self.ot.init_process()

    def _extract_instance_info(self, target):
        """
        extract instance_number, instance_name
        """
        query = PptSql.SELECT_INSTANCE_NUMBER
        for df in self.ot.get_data_by_query(query):
            value_list = df[target].unique().tolist()
            value_list.sort(reverse=False)
            return value_list

    def _extract_period(self):
        s_date = self.config["args"]["s_date"]
        interval = self.config["args"]["interval"]
        s_time = self.config["args"]["s_time"]
        s_interval = self.config["args"]["s_interval"]

        start_date, end_date, start_time, end_time = DateUtils.get_each_date_by_interval2(
            s_date, interval, s_time, s_interval, "%Y.%m.%d"
        )

        period = f"■ 추출기간 : {start_date} ~ {end_date}"
        return period

    def _extract_metric_name_list(self):
        """
        extract metric_name
        """
        metric_name_list = [
            "Host CPU Utilization (%)",
            "Average Active Sessions",
            "Executions Per Sec",
            "Logical Reads Per Sec",
            "Physical Reads Per Sec",
            "Hard Parse Count Per Sec",
        ]

        config_report = Config("report").get_config()
        metric_name_list.extend(config_report["sys_metric"])

        return metric_name_list

    def _extract_metric_finding_text(self, metric_name):
        metric_dict = {
            "Host CPU Utilization (%)": "CPU Usage는 OS 커널에서 가져온 값으로, 전체 시스템에서 사용하는 CPU를 의미.\n"
            "CPU Usage = SYS CPU + User CPU + IO Wait.\n"
            "RAC의 경우, Fail over를 대비하여 각 인스턴스의 CPU 사용률을 50% 이하로 유지되도록 권장.",
            "Average Active Sessions": "Average Active Session은 백그라운드 프로세스를 제외한 User session의 총합을 의미.\n"
            "Oracle의 Stat과 Wait이 모두 반영된 지표이며, 대기 이벤트와 비례하여 증가하기 때문에 시스템의 이상 징후를 인지할 때 가장 중요함.",
            "Executions Per Sec": "SQL Execution은 세션이 수행하는 SQL 수행 횟수의 총합.\n"
            "Execution count가 급감할 경우, AP와 DB 연계 및 대기 이벤트 발생 여부 검토 필요.",
            "Logical Reads Per Sec": "Logical Reads는 세션이 buffer cache에서 읽어 들인 block 총수의 합.\n"
            "Database의 실제 일 량을 가장 직관적으로 판단할 수 있는 지표.",
            "Physical Reads Per Sec": "Physical Reads는 disk에서 읽어 들인 block 총수의 합.\n"
            "OLTP에서 적정치는 logical reads의 1/10~1/100 이하이며, 이 수치보다 클 경우는 SQL 튜닝 필요.",
            "Hard Parse Count Per Sec": "Hard Parse는 Shared pool의 Library cache에서 동일한 SQL을 Hit 하지 못하여 새로 파싱한 횟수를 의미.\n"
            "수치가 높을수록 Shared pool을 비효율적으로 사용하며, literal SQL에 대한 모니터링 및 조치 필요.",
        }

        if metric_name in metric_dict.keys():
            return metric_dict[metric_name]

        else:
            return ""

    def _convert_sql_to_df(self, sql_path, filename, event_name="", inst_num="", schema_name=""):
        """
        sql query문을 dataframe 형태로 변환
        """
        metric_name_list = self._extract_metric_name_list()
        unpack_metric_name_list = str(metric_name_list)[1:-1]

        s_date, e_date, s_time, e_time = DateUtils.get_each_date_by_interval2(
            self.config["args"]["s_date"],
            int(self.config["args"]["interval"]) + 1,
            self.config["args"]["s_time"],
            self.config["args"]["s_interval"],
            arg_fmt="%Y%m%d",
        )

        date_dict = {
            "StartDate": s_date,
            "EndDate": e_date,
            "StartTime": s_time,
            "EndTime": e_time,
            "Metric_Name": unpack_metric_name_list,
            "EVENT_NAME": event_name,
            "INST_NUM": inst_num,
            "SCHEMA_NAME": schema_name,
        }

        query = SystemUtils.get_file_content_in_path(sql_path, filename + ".txt")
        date_query = SqlUtils.sql_replace_to_dict(query, date_dict)

        for df in self.ot.get_data_by_query(date_query):
            df = df.fillna("")

        return df

    def is_even(self, i):
        """
        짝수인 것 반환
        """
        return i % 2 == 0

    def _convert_query_to_df(self, query):
        for df in self.ot.get_data_by_query(query):
            return df

    def _execute_first_slide(self):
        first_slide = self.presentation.slides[0]

        SlideManager.create_text_box(
            first_slide,
            self.position["first_slide"]["date_position"],
            self.today_date,
            Pt(11),
            "나눔스퀘어",
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.WHILE_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

    def _execute_activity_info(self):
        self.logger.info("activity_info.pptx")

        query = PptSql.SELECT_ACTIVITY_INFO
        df = self._convert_query_to_df(query)
        val = df.values[0].astype(str)[0]

        table_row_count = 5
        table_col_count = 2

        tp = SlideManager.convert_inches_to_data(self.position["activity_info"]["table_position"])
        column_width_inches = [Inches(i) for i in self.position["activity_info"]["column_width_inches"]]
        row_height_inches = [Inches(i) for i in self.position["activity_info"]["row_height_inches"]]
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.ACTIVITY_INFO)

        SlideManager.create_text_box(
            target_slide,
            self.position["top_menu"]["left_top_text_position"],
            "Activity Info",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.add_table(target_slide, table_row_count, table_col_count, tp, Pt(11), Pt(9))
        SlideManager.adjust_row_column_size(target_slide, column_width_inches, row_height_inches)
        SlideManager.activity_info_style(target_slide, val)

    def _execute_schedule(self):
        self.logger.info("schedule.pptx")

        tp = SlideManager.convert_inches_to_data(self.position["schedule_info"]["table_position"])
        column_width_inches = [Inches(i) for i in self.position["schedule_info"]["column_width_inches"]]
        row_height_inches = [Inches(i) for i in self.position["schedule_info"]["row_height_inches"]]
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.SCHEDULE)

        table_row_count = 6
        table_col_count = 4

        SlideManager.create_text_box(
            target_slide,
            self.position["top_menu"]["left_top_text_position"],
            "Schedule",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.add_table(target_slide, table_row_count, table_col_count, tp, Pt(11), Pt(9))
        SlideManager.adjust_row_column_size(target_slide, column_width_inches, row_height_inches)
        SlideManager.schedule_style(target_slide)

    def _execute_load_profile(self):
        self.logger.info("load_profile.pptx")
        load_profile_df = self._convert_sql_to_df(self.sql_path, "LOAD_PROFILE")

        tp = SlideManager.convert_inches_to_data(self.position["load_profile"]["table_position"])
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.LOAD_PROFILE)
        row_height = [Inches(i) for i in self.position["load_profile"]["row_height_inches"]]

        SlideManager.create_text_box(
            target_slide,
            self.position["top_menu"]["left_top_text_position"],
            "Load Profile",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        for i in range(len(self.instance_number)):
            if i == 0:
                self._load_profile_detail(load_profile_df, self.instance_name[i : i + 2], target_slide, tp, row_height)

            if i > 1 and i % 2 == 0:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)

                SlideManager.create_text_box(
                    slide_to_add,
                    self.position["top_menu"]["left_top_text_position"],
                    "Load Profile",
                    PerformanceAnalyzer.LARGE_FONT_SIZE,
                    PerformanceAnalyzer.FONT_NAME,
                    PerformanceAnalyzer.BOLD_TRUE,
                    PerformanceAnalyzer.BLACK_COLOR,
                    PerformanceAnalyzer.BASE_LINE_SPACE,
                )
                self._load_profile_detail(load_profile_df, self.instance_name[i : i + 2], slide_to_add, tp, row_height)

    def _load_profile_detail(self, df, inst_name_range, slide, tp, row_height):
        """load_profile_detail"""
        inst_name_df = df[df["sid"].isin(inst_name_range)]
        arranged_df = self._arrange_load_profile(inst_name_df)
        SlideManager.make_table(arranged_df, tp, slide, Pt(11), Pt(9))
        SlideManager.font_center_alignment(slide)

        if len(inst_name_range) == 1:
            SlideManager.adjust_row_column_size(slide, row_height_inches=row_height)

        if len(inst_name_range) != 1:
            SlideManager.load_profile_style(slide)

    def _arrange_load_profile(self, load_profile_df):
        pivot_df = load_profile_df.pivot_table(
            index=["metric_name", "unit", "sid"], columns="mmdd", values="value", aggfunc="first"
        )

        metric_name_mapping = {
            "Host CPU Utilization (%)": "CPU Usage",
            "Average Active Sessions": "Average Active Sessions",
            "User Calls Per Sec": "User Calls",
            "Executions Per Sec": "SQL Executions",
            "User Transaction Per Sec": "TPS",
            "Logical Reads Per Sec": "Logical Reads",
            "Physical Reads Per Sec": "Physical Reads",
            "Hard Parse Count Per Sec": "Hard Parse",
        }

        arranged_df = pivot_df.copy()
        arranged_df.index = arranged_df.index.set_levels(
            arranged_df.index.levels[0].map(metric_name_mapping), level="metric_name"
        )

        metric_name_index = [
            "CPU Usage",
            "Average Active Sessions",
            "User Calls",
            "SQL Executions",
            "TPS",
            "Logical Reads",
            "Physical Reads",
            "Hard Parse",
        ]

        df_sorted_custom = arranged_df.sort_values(
            by=["metric_name"], key=lambda x: pd.Categorical(x, categories=metric_name_index, ordered=True)
        )

        df_sorted_custom = df_sorted_custom.rename_axis(index={"metric_name": "항목", "unit": "Unit", "sid": "SID"})
        test_df = df_sorted_custom.reset_index()

        duplicated_df = test_df.duplicated(["항목", "Unit"], keep="first")
        test_df.loc[duplicated_df, ["항목", "Unit"]] = ""

        return test_df

    def _execute_time_model(self):
        """
        time_model query에 해당하는 레포트
        """
        self.logger.info("time_model.pptx")

        left_tp = SlideManager.convert_inches_to_data(self.position["time_model"]["left_table"])
        right_tp = SlideManager.convert_inches_to_data(self.position["time_model"]["right_table"])
        column_width_inches = [Inches(i) for i in self.position["time_model"]["column_width_inches"]]
        time_model_df = self._convert_sql_to_df(self.sql_path, "TIME_MODEL")

        time_model_df = time_model_df.rename(
            columns={
                "instance_number": "INSTANCE_NUMBER",
                "stat_name": "Statistic Name",
                "time_s": "Time (s)",
                "of_db_time": "% of DB Time",
                "of_total_cpu_time": "% of Total CPU Time",
            }
        )

        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.TIME_MODEL)

        SlideManager.create_text_box(
            target_slide,
            self.position["top_menu"]["left_top_text_position"],
            "Time Model Statistics",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        for i in range(len(self.instance_number)):
            if i == 0:
                self._time_model_detail(
                    time_model_df,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    target_slide,
                    left_tp,
                    right_tp,
                    column_width_inches,
                )

            if i > 1 and i % 2 == 0:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                SlideManager.create_text_box(
                    slide_to_add,
                    self.position["top_menu"]["left_top_text_position"],
                    "Time Model Statistics",
                    PerformanceAnalyzer.LARGE_FONT_SIZE,
                    PerformanceAnalyzer.FONT_NAME,
                    PerformanceAnalyzer.BOLD_TRUE,
                    PerformanceAnalyzer.BLACK_COLOR,
                    PerformanceAnalyzer.BASE_LINE_SPACE,
                )

                self._time_model_detail(
                    time_model_df,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    slide_to_add,
                    left_tp,
                    right_tp,
                    column_width_inches,
                )

    def _time_model_detail(
        self,
        df,
        instance_num_range,
        instance_name_range,
        slide,
        left_table_position,
        right_table_position,
        column_width_inches,
    ):
        """time_model 공통 부분"""
        for i in range(len(instance_num_range)):
            inst_df = df[df["INSTANCE_NUMBER"] == instance_num_range[i]]
            inst_num_drop_df = inst_df.drop(columns=["INSTANCE_NUMBER"])

            rec_position = (
                self.position["time_model"]["left_rectangle_position"]
                if self.is_even(i)
                else self.position["time_model"]["right_rectangle_position"]
            )

            instance_name_position = (
                self.position["time_model"]["left_position"]
                if self.is_even(i)
                else self.position["time_model"]["right_position"]
            )

            tp = left_table_position if self.is_even(i) else right_table_position

            SlideManager.create_shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                rec_position,
                PerformanceAnalyzer.YELLOW_COLOR,
                PerformanceAnalyzer.YELLOW_COLOR,
                0,
            )

            SlideManager.create_text_box(
                slide,
                instance_name_position,
                instance_name_range[i],
                PerformanceAnalyzer.DEFAULT_FONT_SIZE,
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_TRUE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            SlideManager.make_table(inst_num_drop_df, tp, slide, Pt(8), Pt(9), column_width_inches)
            SlideManager.time_model_table_style(slide, RGBColor(242, 242, 242))

    def _execute_db_system(self):
        """
        db_system에 해당하는 레포트
        """
        self.logger.info("db_system.pptx")

        column_width_inches = [Inches(i) for i in self.position["db_system"]["column_width_inches"]]
        tp = SlideManager.convert_inches_to_data(self.position["db_system"]["table_position"])
        df = self._convert_sql_to_df(self.sql_path, "DB_SYSTEM")
        df = df.rename(columns={"category": "구분"})

        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.DB_SYSTEM)

        for i in range(len(self.instance_number)):
            if i == 0:
                self._db_system_detail(
                    df,
                    self.instance_number[i : i + 2],
                    target_slide,
                    column_width_inches,
                    tp,
                    self.instance_name[i : i + 2],
                )

            if i > 1 and i % 2 == 0:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._db_system_detail(
                    df,
                    self.instance_number[i : i + 2],
                    slide_to_add,
                    column_width_inches,
                    tp,
                    self.instance_name[i : i + 2],
                )

    def _db_system_detail(self, df, instance_num_range, slide, column_width_inches, tp, inst_name_range):
        """db_system 공통 부분"""
        list_df = []
        for inst_num in instance_num_range:
            inst_df = df[df["instance_number"] == inst_num]
            list_df.append(inst_df)

        preprocessed_df = self._arrange_db_system(list_df, inst_name_range)

        SlideManager.create_text_box(
            slide,
            self.position["top_menu"]["left_top_text_position"],
            "DB SYSTEM",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )
        SlideManager.make_table(preprocessed_df, tp, slide, Pt(11), Pt(9), column_width_inches)
        SlideManager.db_system_table_style(slide)

    def _arrange_db_system(self, list_df, inst_name_range):
        """db_system dataframe 전처리"""
        dfs_to_merge = []
        new_row = pd.DataFrame({"구분": "Business"}, index=[0])

        for idx, df in enumerate(list_df):
            inst_num_drop_df = df.drop(columns=["instance_number"])
            concat_df = pd.concat([new_row, inst_num_drop_df]).reset_index(drop=True)
            inst_name = inst_name_range[idx]
            concat_df.rename(columns={"value": f"{inst_name}"}, inplace=True)
            arranged_df = concat_df.replace(
                [
                    "HOST_NAME",
                    "INSTANCE_NAME",
                    "PLATFORM_NAME",
                    "BANNER",
                    "NUM_CPUS",
                    "PHYSICAL_MEMORY_BYTES",
                    "Buffer Cache Size",
                    "Shared Pool Size",
                    "Large Pool Size",
                    "Java Pool Size",
                    "Streams Pool Size",
                    "pga_aggregate_limit",
                ],
                [
                    "Hostname",
                    "Instance Name",
                    "OS Version",
                    "GI & DB Version",
                    "OS CPU",
                    "OS Memory",
                    "Buffer Cache",
                    "Shared Pool",
                    "Large Pool",
                    "Java Pool",
                    "Streams Pool",
                    "PGA",
                ],
            )

            arranged_df = arranged_df.set_index("구분")
            arranged_df = arranged_df.reindex(
                [
                    "Business",
                    "Hostname",
                    "OS Version",
                    "GI & DB Version",
                    "Instance Name",
                    "OS CPU",
                    "OS Memory",
                    "Buffer Cache",
                    "Shared Pool",
                    "Java Pool",
                    "Large Pool",
                    "Streams Pool",
                    "PGA",
                ]
            ).reset_index(drop=False)

            arranged_df = arranged_df.fillna("")
            dfs_to_merge.append(arranged_df)

        if len(dfs_to_merge) == 1:
            arranged_df.insert(loc=0, column="", value="")
            return arranged_df

        else:
            merge_df = pd.merge(dfs_to_merge[0], dfs_to_merge[1], on="구분", how="outer")
            merge_df.insert(loc=0, column="", value="")
            return merge_df

    def _execute_metric(self):
        """
        metric에 해당하는 레포트
        """
        self.logger.info("metric.pptx")
        self._convert_sql_to_df(self.sql_path, "METRIC")
        df = self._convert_sql_to_df(self.sql_path, "METRIC")

        metric_name_list = self._extract_metric_name_list()
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.METRIC)

        for i in range(len(metric_name_list)):
            if i == 0:
                green_text = self._extract_metric_finding_text(metric_name_list[i])
                self._metric_detail(target_slide, metric_name_list[i], df, green_text)

            else:
                green_text = self._extract_metric_finding_text(metric_name_list[i])
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._metric_detail(slide_to_add, metric_name_list[i], df, green_text)

    def _metric_detail(self, slide, metric_name, df, green_text):
        """metric 공통 부분"""
        new_metric_name = (
            PerformanceAnalyzer.metric_top_menu_dict[metric_name]
            if metric_name in PerformanceAnalyzer.metric_top_menu_dict
            else metric_name
        )

        SlideManager.create_text_box(
            slide,
            self.position["top_menu"]["left_top_text_position"],
            new_metric_name,
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            self.position["common_use"]["extract_period"],
            self.period,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["metric"]["avg_rectangle_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            slide,
            self.position["metric"]["avg_text_position"],
            "AVG",
            PerformanceAnalyzer.MEDIUM_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["metric"]["max_rectangle_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            slide,
            self.position["metric"]["max_text_position"],
            "MAX",
            PerformanceAnalyzer.MEDIUM_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        y_axis_val = self._extract_y_axis_value(metric_name)

        SlideManager.create_text_box(
            slide,
            self.position["y_axis"]["upper_position"],
            y_axis_val,
            Pt(9),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            self.position["y_axis"]["down_position"],
            y_axis_val,
            Pt(9),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        metric_df = df[df["metric_name"] == metric_name]
        copied_df = metric_df.copy()
        preprocessed_metric_df = self._set_df_date_time(copied_df)

        ag_y_axis_max_value, ag_value_tuple_list = self._extract_df_value(preprocessed_metric_df, "ag")
        mx_y_axis_max_value, mx_value_tuple_list = self._extract_df_value(preprocessed_metric_df, "mx")

        ag_chart = SlideManager.insert_line_chart(
            slide,
            preprocessed_metric_df,
            self.position["metric"]["avg_coordinate"],
            self.instance_number,
            self.instance_name,
            "ag",
        )

        mx_chart = SlideManager.insert_line_chart(
            slide,
            preprocessed_metric_df,
            self.position["metric"]["max_coordinate"],
            self.instance_number,
            self.instance_name,
            "mx",
        )

        self._chart_style(slide, ag_value_tuple_list, mx_value_tuple_list)
        SlideManager.set_y_axis_max_value(ag_chart, mx_y_axis_max_value)
        SlideManager.set_y_axis_max_value(mx_chart, mx_y_axis_max_value)

        bottom_text = self._extract_chart_bottom_text(
            slide, metric_name, preprocessed_metric_df, "ag", "mx", "AVG", "MAX"
        )

        self._set_bottom_exclude_text(
            slide,
            self.position["metric"]["straight_position"],
            self.position["metric"]["right_triangle_position"],
            self.position["metric"]["bottom_findings_position"],
        )

        self._insert_memory_bottom_text(slide, self.position["metric"]["bottom_text_position"], green_text, bottom_text)

    def _extract_df_value(self, df, col):
        max_list = []
        value_list = []

        for idx, inst_num in enumerate(self.instance_number):
            extract_df = df[df["instance_number"] == inst_num]

            max_score = extract_df[col].max()
            max_list.append(max_score)

            col_value_tuple = tuple(extract_df[col].tolist())
            k_value_tuple = SlideManager.convert_unit(col_value_tuple)
            value_list.append(k_value_tuple)

        y_axis_max_value = SlideManager.make_max_value(max_list)

        return y_axis_max_value, value_list

    def _execute_rac(self):
        """
        rac에 해당하는 레포트
        """
        self.logger.info("RAC.pptx")

        if len(self.instance_name) != 1:
            df = self._convert_sql_to_df(self.sql_path, "RAC")
            preprocessed_df = self._set_df_date_time(df)
            _, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.RAC)

            SlideManager.create_text_box(
                target_slide,
                self.position["top_menu"]["left_top_text_position"],
                "RAC Interconnect traffic",
                PerformanceAnalyzer.LARGE_FONT_SIZE,
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_TRUE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            SlideManager.create_text_box(
                target_slide,
                self.position["common_use"]["extract_period"],
                self.period,
                PerformanceAnalyzer.SMALL_FONT_SIZE,
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_FALSE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            SlideManager.create_shape(
                target_slide,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                self.position["rac"]["rec1_position"],
                PerformanceAnalyzer.YELLOW_COLOR,
                PerformanceAnalyzer.YELLOW_COLOR,
                0,
            )

            SlideManager.create_text_box(
                target_slide,
                self.position["rac"]["sent_position"],
                "DB Block Sent",
                PerformanceAnalyzer.MEDIUM_FONT_SIZE,
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_TRUE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            SlideManager.create_shape(
                target_slide,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                self.position["rac"]["rec2_position"],
                PerformanceAnalyzer.YELLOW_COLOR,
                PerformanceAnalyzer.YELLOW_COLOR,
                0,
            )

            SlideManager.create_text_box(
                target_slide,
                self.position["rac"]["received_position"],
                "DB Block Received",
                PerformanceAnalyzer.MEDIUM_FONT_SIZE,
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_TRUE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            SlideManager.create_text_box(
                target_slide,
                self.position["y_axis"]["upper_position"],
                "Block_sent(kb/sec)",
                Pt(9),
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_FALSE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            SlideManager.create_text_box(
                target_slide,
                self.position["y_axis"]["down_position"],
                "Block received(kb/sec)",
                Pt(9),
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_FALSE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            block_sent_y_axis_max_value, block_value_tuple_list = self._extract_df_value(
                preprocessed_df, "bytes_sentpsec"
            )
            block_received_y_axis_max_value, block_received_value_tuple_list = self._extract_df_value(
                preprocessed_df, "bytes_receivedpsec"
            )

            block_sent_chart = SlideManager.insert_line_chart(
                target_slide,
                preprocessed_df,
                self.position["rac"]["sent_coordinate"],
                self.instance_number,
                self.instance_name,
                "bytes_sentpsec",
            )

            block_received_chart = SlideManager.insert_line_chart(
                target_slide,
                preprocessed_df,
                self.position["rac"]["received_coordinate"],
                self.instance_number,
                self.instance_name,
                "bytes_receivedpsec",
            )

            self._chart_style(target_slide, block_value_tuple_list, block_received_value_tuple_list)

            SlideManager.set_y_axis_max_value(block_sent_chart, block_sent_y_axis_max_value)
            SlideManager.set_y_axis_max_value(block_received_chart, block_received_y_axis_max_value)

            green_text = (
                "RAC Interconnect Traffic은 RAC 인스턴스간에 Cache fusion으로 주고받은 블록 사이즈를 의미.\n"
                "Interconnect trafic으로 인한 병목이 발생하면, 관련 인스턴스들의 I/O가 지연됨에 따라, 시스템 성능이 저하될 수 있음. 해당시간에 수행되는 SQL 튜닝 필요."
            )

            bottom_text = self._extract_chart_bottom_text(
                target_slide,
                "RAC Interconnect traffic",
                preprocessed_df,
                "bytes_sentpsec",
                "bytes_receivedpsec",
                "DB Block Sent",
                "DB Block Received",
            )

            self._set_bottom_exclude_text(
                target_slide,
                self.position["metric"]["straight_position"],
                self.position["metric"]["right_triangle_position"],
                self.position["metric"]["bottom_findings_position"],
            )

            self._insert_memory_bottom_text(
                target_slide, self.position["metric"]["bottom_text_position"], green_text, bottom_text
            )

        else:
            pass

    def _execute_top_n_wait_events(self):
        """
        top_n_wait_events에 해당하는 레포트
        """
        self.logger.info("top_n_wait_events.pptx")

        left_tp = SlideManager.convert_inches_to_data(self.position["top_n"]["left_table"])
        right_tp = SlideManager.convert_inches_to_data(self.position["top_n"]["right_table"])
        column_width_inches = [Inches(i) for i in self.position["top_n"]["column_width_inches"]]

        top_n_df = self._convert_sql_to_df(self.sql_path, "TOP_N_Wait_Events")
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.TON_N)

        for i in range(len(self.instance_number)):
            if i == 0:
                SlideManager.create_text_box(
                    target_slide,
                    self.position["top_menu"]["left_top_text_position"],
                    "TOP-N Wait Events",
                    PerformanceAnalyzer.LARGE_FONT_SIZE,
                    PerformanceAnalyzer.FONT_NAME,
                    PerformanceAnalyzer.BOLD_TRUE,
                    PerformanceAnalyzer.BLACK_COLOR,
                    PerformanceAnalyzer.BASE_LINE_SPACE,
                )

                self._top_n_wait_events_detail(
                    target_slide,
                    top_n_df,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    left_tp,
                    right_tp,
                    column_width_inches,
                )

            if i > 1 and i % 2 == 0:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                SlideManager.create_text_box(
                    slide_to_add,
                    self.position["top_menu"]["left_top_text_position"],
                    "TOP-N Wait Events",
                    PerformanceAnalyzer.LARGE_FONT_SIZE,
                    PerformanceAnalyzer.FONT_NAME,
                    PerformanceAnalyzer.BOLD_TRUE,
                    PerformanceAnalyzer.BLACK_COLOR,
                    PerformanceAnalyzer.BASE_LINE_SPACE,
                )

                self._top_n_wait_events_detail(
                    slide_to_add,
                    top_n_df,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    left_tp,
                    right_tp,
                    column_width_inches,
                )

    def _top_n_wait_events_detail(
        self, slide, df, instance_num_range, instance_name_range, left_tp, right_tp, column_width_inches
    ):
        """top_n 공통 부분"""
        green_text = "DB CPU Time을 포함하여 Top Event 발생 현황 출력."
        inst_event_dict = self._extract_top3_list(df)
        bottom_text = self._extract_top_n_bottom_text(inst_event_dict, instance_num_range, instance_name_range)

        for i in range(len(instance_num_range)):
            inst_df = df[df["instance_number"] == instance_num_range[i]]
            inst_num_drop_df = inst_df.drop(columns=["instance_number"])
            inst_num_drop_df.rename(
                columns={
                    "event_name": "Event",
                    "wait_class": "Wait class",
                    "total_waits": "Waits",
                    "total_wait_time_sec": "Time (s)",
                    "avg_wait": "Avg wait (ms)",
                    "db_time_ratio": "% DB time",
                },
                inplace=True,
            )
            rec_position = (
                self.position["top_n"]["left_rectangle_position"]
                if self.is_even(i)
                else self.position["top_n"]["right_rectangle_position"]
            )

            instance_name_position = (
                self.position["top_n"]["left_position"] if self.is_even(i) else self.position["top_n"]["right_position"]
            )

            tp = left_tp if self.is_even(i) else right_tp

            SlideManager.create_shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                rec_position,
                PerformanceAnalyzer.YELLOW_COLOR,
                PerformanceAnalyzer.YELLOW_COLOR,
                0,
            )

            SlideManager.create_text_box(
                slide,
                instance_name_position,
                instance_name_range[i],
                PerformanceAnalyzer.DEFAULT_FONT_SIZE,
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_TRUE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            SlideManager.make_table(
                inst_num_drop_df,
                tp,
                slide,
                PerformanceAnalyzer.SMALL_FONT_SIZE,
                PerformanceAnalyzer.SMALL_FONT_SIZE,
                column_width_inches,
            )

            self._set_bottom_exclude_text(
                slide,
                self.position["metric"]["straight_position"],
                self.position["metric"]["right_triangle_position"],
                self.position["metric"]["bottom_findings_position"],
            )

            self._insert_memory_bottom_text(
                slide, self.position["metric"]["bottom_text_position"], green_text, bottom_text
            )

    def _execute_top_3_wait_events(self):
        """
        top_n_wait_events에서 뽑은 top3에 해당하는 레포트
        """
        self.logger.info("top_3_wait_events.pptx")
        df = self._convert_sql_to_df(self.sql_path, "TOP_N_Wait_Events")
        instance_dict = self._extract_top3_list(df)

        upper_tp = SlideManager.convert_inches_to_data(self.position["top_3"]["upper_table"])
        down_tp = SlideManager.convert_inches_to_data(self.position["top_3"]["down_table"])
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.TOP_3)

        k = 0

        for indx, (key, value) in enumerate(instance_dict.items()):
            inst_num = key.split("_")[-1]
            for event_name in value:
                sql1_df = self._extract_preprocessed_df(event_name, inst_num, "TOP_1_Wait_Events")

                sql2_df = self._extract_preprocessed_df(event_name, inst_num, "TOP_2_Wait_Events", ["instance_number"])
                sql2_df.rename(
                    columns={
                        "time_outs": "%Time -outs",
                        "total_wait_time": "Total Wait Time (s)",
                        "avg_wait": "Avg wait (ms)",
                        "waits_txn": "Waits /txn",
                    }
                )

                sql3_preprocess_df = self._extract_preprocessed_df(
                    event_name, inst_num, "TOP_3_Wait_Events", ["instance_number"]
                )

                sql3_preprocess_df.rename(
                    columns={
                        "sql_id": "SQL_ID",
                        "sql_plan_hash_value": "Plan Hash Value",
                        "event": "Event",
                        "sql_opname": "TOP ROW SOURCE",
                        "sql_text": "SQL Text",
                    }
                )

                sql3_preprocess_df = self._set_blank_df(sql3_preprocess_df.head(3), 3)
                sql3_df = self._adjust_df_length(sql3_preprocess_df, "sql_text")

                if k == 0:
                    self._top3_detail(
                        target_slide,
                        event_name,
                        self.instance_number[indx],
                        self.instance_name[indx],
                        sql1_df,
                        sql2_df,
                        sql3_df,
                        upper_tp,
                        down_tp,
                    )
                    k += 1

                elif k >= 1:
                    slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                    self._top3_detail(
                        slide_to_add,
                        event_name,
                        self.instance_number[indx],
                        self.instance_name[indx],
                        sql1_df,
                        sql2_df,
                        sql3_df,
                        upper_tp,
                        down_tp,
                    )
                    k += 1

    def _top3_detail(self, slide, event_name, inst_num, inst_name, sql1_df, sql2_df, sql3_df, upper_tp, down_tp):
        """top3 공통 부분"""
        SlideManager.create_text_box(
            slide,
            self.position["top_menu"]["left_top_text_position"],
            event_name,
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            self.position["top_3"]["extract_period"],
            self.period,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["top_3"]["ora_rectangle_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            slide,
            self.position["top_3"]["ora_position"],
            inst_name,
            PerformanceAnalyzer.DEFAULT_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )
        upper_column_width_inches = [Inches(i) for i in self.position["top_3"]["upper_column_width_inches"]]
        down_column_width_inches = [Inches(i) for i in self.position["top_3"]["down_column_width_inches"]]
        bottom_text = self._extract_top3_text(inst_name, event_name, sql1_df, sql2_df)

        line_chart = SlideManager.insert_chart(
            slide,
            sql1_df["combined_datetime"].tolist(),
            tuple(sql1_df["mspsec"].tolist()),
            inst_name,
            self.position["top_3"]["chart_position"],
        )

        y_axis_max_value = SlideManager.make_max_value([sql1_df["mspsec"].max()])
        self._top3_chart_detail(slide, inst_num, sql1_df["mspsec"])
        SlideManager.set_y_axis_max_value(line_chart, y_axis_max_value)
        SlideManager.make_table(
            sql2_df,
            upper_tp,
            slide,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            upper_column_width_inches,
        )
        SlideManager.make_table(
            sql3_df,
            down_tp,
            slide,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            down_column_width_inches,
        )

        self._set_bottom_text(
            slide,
            bottom_text,
            self.position["metric"]["straight_position"],
            self.position["metric"]["right_triangle_position"],
            self.position["metric"]["bottom_findings_position"],
            self.position["metric"]["bottom_text_position"],
        )

    def _extract_name_value_list(self):
        """memory blue_box에 들어가는 text 정리"""
        amm_memory_df = self._convert_sql_to_df(self.sql_path, "MEMORY_AMM")

        amm_memory_list = ["Physical memory", "sga_target", "db_block_size", "pga_aggregate_target"]
        name_value_list = []
        for name in amm_memory_list:
            filtered_row = amm_memory_df[amm_memory_df["NAME"] == name]
            amm_memory_value = filtered_row["VALUE"].iloc[0]

            if name == "Physical memory":
                changed_name = "Physical Memory"

            else:
                changed_name = name.upper()

            make_sentence = f"{changed_name} = {amm_memory_value}"
            name_value_list.append(make_sentence)

        return name_value_list

    def _execute_memory_advice_sga(self):
        self.logger.info("memory_advice_sga.pptx")

        left_tp = self.position["memory_advice"]["left_table"]
        right_tp = self.position["memory_advice"]["right_table"]

        df = self._convert_sql_to_df(self.sql_path, "MEMORY_SGA")
        df = df.rename(columns={"physical_reads": "PHYSICAL_READS", "db_time": "DB_TIME"})

        finding_df = self._convert_sql_to_df(self.sql_path, "MEMORY_SGA_SQL2")
        finding_df_col1 = "sga_size"
        finding_df_col2 = "physical_reads"

        sga_color_list = [RGBColor(32, 56, 100), RGBColor(143, 170, 220)]
        top_menu_text = "Memory Advice – SGA"
        category_col = "sga_size"
        series_col1 = "PHYSICAL_READS"
        series_col2 = "DB_TIME"
        bottom_text1 = "SGA"
        bottom_text2 = "DISK I/O"
        bottom_green_text = "sga_target 파라미터 설정시(0 아닌 값) 차트 확인 가능."

        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.MEMORY_SGA)

        for i in range(len(self.instance_number)):
            if i == 0:
                self._memory_advice_detail(
                    target_slide,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    top_menu_text,
                    bottom_text1,
                    bottom_text2,
                    bottom_green_text,
                    finding_df,
                    finding_df_col1,
                    finding_df_col2,
                )
                self._memory_advice_chart(
                    target_slide,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    left_tp,
                    right_tp,
                    sga_color_list,
                    df,
                    category_col,
                    series_col1,
                    series_col2,
                )

            elif i > 1 and i % 2 == 0:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._memory_advice_detail(
                    slide_to_add,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    top_menu_text,
                    bottom_text1,
                    bottom_text2,
                    bottom_green_text,
                    finding_df,
                    finding_df_col1,
                    finding_df_col2,
                )
                self._memory_advice_chart(
                    slide_to_add,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    left_tp,
                    right_tp,
                    sga_color_list,
                    df,
                    category_col,
                    series_col1,
                    series_col2,
                )

    def _execute_memory_advice_pga(self):
        self.logger.info("memory_advice_pga.pptx")

        left_tp = self.position["memory_advice"]["left_table"]
        right_tp = self.position["memory_advice"]["right_table"]

        df = self._convert_sql_to_df(self.sql_path, "MEMORY_PGA")
        df = df.rename(columns={"pga_cache_hit_ratio": "PGA_CACHE_HIT", "overalloc_ratio": "OVERALLOC"})

        finding_df = self._convert_sql_to_df(self.sql_path, "MEMORY_PGA_SQL2")
        finding_df_col1 = "pga_size"
        finding_df_col2 = "pga_cache_hit"
        pga_color_list = [RGBColor(127, 96, 0), RGBColor(255, 217, 102)]
        top_menu_text = "Memory Advice – PGA"
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.MEMORY_PGA)
        category_col = "pga_size"
        series_col1 = "PGA_CACHE_HIT"
        series_col2 = "OVERALLOC"
        bottom_text1 = "PGA"
        bottom_text2 = "PGA_CACHE_HIT"
        bottom_green_text = "pga_aggregate_target 파라미터 설정시(0 아닌 값) 차트 확인 가능."

        for i in range(len(self.instance_number)):
            if i == 0:
                self._memory_advice_detail(
                    target_slide,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    top_menu_text,
                    bottom_text1,
                    bottom_text2,
                    bottom_green_text,
                    finding_df,
                    finding_df_col1,
                    finding_df_col2,
                )
                self._memory_advice_chart(
                    target_slide,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    left_tp,
                    right_tp,
                    pga_color_list,
                    df,
                    category_col,
                    series_col1,
                    series_col2,
                )

            elif i > 1 and i % 2 == 0:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._memory_advice_detail(
                    slide_to_add,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    top_menu_text,
                    bottom_text1,
                    bottom_text2,
                    bottom_green_text,
                    finding_df,
                    finding_df_col1,
                    finding_df_col2,
                )
                self._memory_advice_chart(
                    slide_to_add,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    left_tp,
                    right_tp,
                    pga_color_list,
                    df,
                    category_col,
                    series_col1,
                    series_col2,
                )

    def _execute_memory_shared_pool(self):
        self.logger.info("memory_advice_shared_pool.pptx")

        left_tp = self.position["memory_advice"]["left_table"]
        right_tp = self.position["memory_advice"]["right_table"]

        df = self._convert_sql_to_df(self.sql_path, "MEMORY_SHARED_POOL")
        df = df.rename(columns={"lc_time_saved": "LC_TIME_SAVED", "lc_memory_object_hit": "LC_MEMORY_OBJECT_HIT"})

        finding_df = self._convert_sql_to_df(self.sql_path, "MEMORY_SHARED_POOL_SQL2")
        finding_df_col1 = "shared_pool_size"
        finding_df_col2 = "lc_time_saved"
        shared_pool_color_list = [RGBColor(56, 87, 35), RGBColor(169, 209, 142)]
        top_menu_text = "Memory Advice – Shared pool"
        slide_num, target_slide = SlideManager.read_slide(
            self.presentation.slides, PerformanceAnalyzer.MEMORY_SHARED_POOL
        )
        category_col = "shared_pool_size"
        series_col1 = "LC_TIME_SAVED"
        series_col2 = "LC_MEMORY_OBJECT_HIT"
        bottom_text1 = "SHARED_POOL"
        bottom_text2 = "LC_TIME_SAVED"
        bottom_green_text = "statistics_level 파라미터 ALL 또는 TYPICAL 설정시 차트 확인 가능."

        for i in range(len(self.instance_number)):
            if i == 0:
                self._memory_advice_detail(
                    target_slide,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    top_menu_text,
                    bottom_text1,
                    bottom_text2,
                    bottom_green_text,
                    finding_df,
                    finding_df_col1,
                    finding_df_col2,
                )
                self._memory_advice_chart(
                    target_slide,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    left_tp,
                    right_tp,
                    shared_pool_color_list,
                    df,
                    category_col,
                    series_col1,
                    series_col2,
                )

            elif i > 1 and i % 2 == 0:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._memory_advice_detail(
                    slide_to_add,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    top_menu_text,
                    bottom_text1,
                    bottom_text2,
                    bottom_green_text,
                    finding_df,
                    finding_df_col1,
                    finding_df_col2,
                )
                self._memory_advice_chart(
                    slide_to_add,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    left_tp,
                    right_tp,
                    shared_pool_color_list,
                    df,
                    category_col,
                    series_col1,
                    series_col2,
                )

    def _execute_memory_db_buffer_cache(self):
        self.logger.info("memory_advice_db_buffer_cache.pptx")

        left_tp = self.position["memory_advice"]["left_table"]
        right_tp = self.position["memory_advice"]["right_table"]

        df = self._convert_sql_to_df(self.sql_path, "MEMORY_DB_BUFFER_CACHE")
        df = df.rename(columns={"physical_reads": "PHYSICAL_READS"})

        finding_df = self._convert_sql_to_df(self.sql_path, "MEMORY_DB_BUFFER_CACHE_SQL2")

        finding_df_col1 = "size_for_estimate"
        finding_df_col2 = "physical_reads"
        db_buffer_cache_color_list = [RGBColor(237, 125, 49)]
        top_menu_text = "Memory Advice – DB buffer cache"
        slide_num, target_slide = SlideManager.read_slide(
            self.presentation.slides, PerformanceAnalyzer.MEMORY_DB_BUFFER_CACHE
        )
        category_col = "size_for_estimate"
        series_col1 = "PHYSICAL_READS"
        bottom_text1 = "DB_BUFFER_CACHE"
        bottom_text2 = "PHYSICAL_READS"
        bottom_green_text = "statistics_level 파라미터 ALL 또는 TYPICAL 설정시 차트 확인 가능."

        for i in range(len(self.instance_number)):
            if i == 0:
                self._memory_advice_detail(
                    target_slide,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    top_menu_text,
                    bottom_text1,
                    bottom_text2,
                    bottom_green_text,
                    finding_df,
                    finding_df_col1,
                    finding_df_col2,
                )
                self._memory_advice_one_bar_chart(
                    target_slide,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    left_tp,
                    right_tp,
                    db_buffer_cache_color_list,
                    df,
                    category_col,
                    series_col1,
                )

            elif i > 1 and i % 2 == 0:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._memory_advice_detail(
                    slide_to_add,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    top_menu_text,
                    bottom_text1,
                    bottom_text2,
                    bottom_green_text,
                    finding_df,
                    finding_df_col1,
                    finding_df_col2,
                )
                self._memory_advice_one_bar_chart(
                    slide_to_add,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    left_tp,
                    right_tp,
                    db_buffer_cache_color_list,
                    df,
                    category_col,
                    series_col1,
                )

    def _memory_advice_one_bar_chart(
        self,
        slide,
        inst_number_range,
        inst_name_range,
        left_tp,
        right_tp,
        memory_color_list,
        df,
        category_col,
        series_col1="",
    ):
        for i in range(len(inst_number_range)):
            rec_position = (
                self.position["memory_advice"]["rectangle_position"]
                if self.is_even(i)
                else self.position["memory_advice"]["rectangle2_position"]
            )

            instance_name_position = (
                self.position["memory_advice"]["text1_position"]
                if self.is_even(i)
                else self.position["memory_advice"]["text2_position"]
            )

            pt = left_tp if self.is_even(i) else right_tp

            inst_df = df[df["inst_id"] == inst_number_range[i]]
            category_df = inst_df[category_col].to_list()

            preprocessed_df = inst_df.copy()
            preprocessed_df[series_col1] = preprocessed_df[series_col1].replace("%", "", regex=True).astype(float)

            df_val1 = preprocessed_df[series_col1]
            df_val1_tuple_list = tuple(df_val1.tolist())

            SlideManager.add_bar_chart(category_df, df_val1_tuple_list, pt, slide, series_col1)
            SlideManager.set_memory_bar_chart_style(slide, memory_color_list)

            SlideManager.create_shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                rec_position,
                PerformanceAnalyzer.YELLOW_COLOR,
                PerformanceAnalyzer.YELLOW_COLOR,
                0,
            )

            SlideManager.create_text_box(
                slide,
                instance_name_position,
                inst_name_range[i],
                Pt(11),
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_TRUE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

    def _memory_advice_detail(
        self,
        slide,
        inst_num_range,
        inst_name_range,
        top_menu_text,
        text1,
        text2,
        bottom_green_text,
        finding_df,
        finding_col1,
        finding_col2,
    ):
        sp = self.position["memory_advice"]["bottom_text_position"]

        SlideManager.create_text_box(
            slide,
            self.position["top_menu"]["left_top_text_position"],
            top_menu_text,
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            self.position["memory_advice"]["select_period"],
            f"■ 조회일자 : {self.today_date}",
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        bottom_text = self._extract_memory_bottom_text(
            inst_num_range, inst_name_range, text1, text2, finding_df, finding_col1, finding_col2
        )
        self._insert_memory_bottom_text(slide, sp, bottom_green_text, bottom_text)

        self._set_bottom_exclude_text(
            slide,
            self.position["memory_advice"]["straight_position"],
            self.position["memory_advice"]["right_triangle_position"],
            self.position["memory_advice"]["bottom_findings_position"],
        )

    def _memory_advice_chart(
        self,
        slide,
        inst_number_range,
        inst_name_range,
        left_tp,
        right_tp,
        memory_color_list,
        df,
        category_col,
        series_col1="",
        series_col2="",
    ):
        for i in range(len(inst_number_range)):
            rec_position = (
                self.position["memory_advice"]["rectangle_position"]
                if self.is_even(i)
                else self.position["memory_advice"]["rectangle2_position"]
            )

            instance_name_position = (
                self.position["memory_advice"]["text1_position"]
                if self.is_even(i)
                else self.position["memory_advice"]["text2_position"]
            )

            pt = left_tp if self.is_even(i) else right_tp

            inst_df = df[df["inst_id"] == inst_number_range[i]]
            category_df = inst_df[category_col].to_list()

            preprocessed_df = inst_df.copy()
            preprocessed_df[series_col1] = preprocessed_df[series_col1].replace("%", "", regex=True).astype(float)
            preprocessed_df[series_col2] = preprocessed_df[series_col2].replace("%", "", regex=True).astype(float)

            df_val1 = preprocessed_df[series_col1]
            df_val2 = preprocessed_df[series_col2]

            SlideManager.add_close_bar_chart(category_df, df_val1, df_val2, series_col1, series_col2, pt, slide)
            SlideManager.set_memory_bar_chart_style(slide, memory_color_list)

            SlideManager.create_shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                rec_position,
                PerformanceAnalyzer.YELLOW_COLOR,
                PerformanceAnalyzer.YELLOW_COLOR,
                0,
            )

            SlideManager.create_text_box(
                slide,
                instance_name_position,
                inst_name_range[i],
                Pt(11),
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_TRUE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

    def _extract_database_parameter(self):
        for i in self.oracle_version:
            if i == "19":
                df = self._convert_sql_to_df(self.sql_path, "DATABASE_PARAMETER_19C_SQL1")
                df = df.rename(
                    columns={
                        "category": "구분",
                        "name": "파라미터",
                        "default_val": "기본값",
                        "value": "현재값",
                        "recommend_val": "권장값",
                        "comm": "설명",
                    }
                )
                text = "Oracle Database 19c"

            elif i == "11":
                df = self._convert_sql_to_df(self.sql_path, "DATABASE_PARAMETER_11G_SQL1")
                df = df.rename(
                    columns={
                        "category": "구분",
                        "name": "파라미터",
                        "default_val": "기본값",
                        "value": "현재값",
                        "recommend_val": "권장값",
                        "comm": "설명",
                    }
                )
                text = "Oracle Database 11g"

            else:
                df = pd.DataFrame()
                text = None

        return df, text

    def _execute_database_parameter(self):
        self.logger.info("database_parameter.pptx")
        df, text = self._extract_database_parameter()

        if not df.empty:
            column_width_inches = [Inches(i) for i in self.position["database_parameter"]["column_width_inches"]]
            table_position = SlideManager.convert_inches_to_data(self.position["database_parameter"]["table_position"])
            rec_position = self.position["database_parameter"]["rectangle_position"]
            text_position = self.position["database_parameter"]["text_position"]

            slide_num, target_slide = SlideManager.read_slide(
                self.presentation.slides, PerformanceAnalyzer.DATABASE_PARAMETER
            )

            SlideManager.create_text_box(
                target_slide,
                self.position["top_menu"]["left_top_text_position"],
                "Database Parameter",
                PerformanceAnalyzer.LARGE_FONT_SIZE,
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_TRUE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            for idx, i in enumerate(range(0, len(df), 10)):
                db_para_sql = df.iloc[i : i + 10].reset_index(drop=True)

                if idx == 0:
                    self._database_parameter_detail(
                        target_slide,
                        db_para_sql,
                        table_position,
                        rec_position,
                        text_position,
                        text,
                        column_width_inches,
                    )

                else:
                    slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                    SlideManager.create_text_box(
                        slide_to_add,
                        self.position["top_menu"]["left_top_text_position"],
                        "Database Parameter",
                        PerformanceAnalyzer.LARGE_FONT_SIZE,
                        PerformanceAnalyzer.FONT_NAME,
                        PerformanceAnalyzer.BOLD_TRUE,
                        PerformanceAnalyzer.BLACK_COLOR,
                        PerformanceAnalyzer.BASE_LINE_SPACE,
                    )
                    self._database_parameter_detail(
                        slide_to_add,
                        db_para_sql,
                        table_position,
                        rec_position,
                        text_position,
                        text,
                        column_width_inches,
                    )

    def _database_parameter_detail(
        self, slide, df, table_position, rec_position, text_position, text, column_width_inches
    ):
        SlideManager.create_text_box(
            slide,
            self.position["common_use"]["select_period"],
            f"■ 조회일자 : {self.today_date}",
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            rec_position,
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            slide,
            text_position,
            text,
            Pt(14),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        duplicated_df = df.duplicated(["구분"], keep="first")
        df.loc[duplicated_df, ["구분"]] = ""

        SlideManager.make_table(df, table_position, slide, Pt(10), Pt(9), column_width_inches)
        SlideManager.database_style(slide, df)

    def _preprocessed_literal_val(self, df):
        df.loc[1:, "SIZE"] = df.loc[1:, "SIZE"].replace("%", "", regex=True).astype(float)
        sum_of_values = df.loc[1:, "SIZE"].sum()
        sum_of_values_rounded = round(sum_of_values, 1)
        df.loc[0, "SIZE"] = 100 - sum_of_values_rounded

        return df

    def top_literal_sql_ratio(self):
        """
        top_literal_sql_ratio
        """
        self.logger.info("top_literal_sql_ratio.pptx")
        slide_num, target_slide = SlideManager.read_slide(
            self.presentation.slides, PerformanceAnalyzer.TOP_LITERAL_RATIO
        )

        for i in range(len(self.instance_number)):
            df = self._convert_sql_to_df(self.sql_path, "TOP_LITERAL_TREE_MAP", inst_num=str(self.instance_number[i]))
            preprocessed_df = self._preprocessed_literal_val(df)

            if i == 0:
                print("00", preprocessed_df)
                self._update_tree_map(target_slide, preprocessed_df, self.instance_name[i])

                # preprocessed_df.sort_values('SIZE', ascending=False, inplace=True)
                # fig, ax = plt.subplots(figsize=(12, 12), dpi=100, subplot_kw=dict(aspect=1.156))
                #
                # # cmap=matplotlib.cm.plasma, matplotlib.cm.viridis
                # tr.treemap(ax, preprocessed_df, area='SIZE', labels='sql_id', cmap='Set2',
                #            rectprops=dict(ec='w'), textprops=dict(c='w'))
                #
                # # squarify.plot(sizes=df['SIZE']*2,
                # #               label=df['sql_id'],
                # #               value=df['SIZE'].astype(str) + '%',
                # #               color=sb.color_palette("Spectral", len(df['SIZE'])),
                # #               alpha=.8,
                # #               text_kwargs={'fontsize': 7})
                #
                # plt.axis('off')
                # plt.show()
                # SlideManager.update_tree_map(target_slide)

            if i >= 1:
                print("11", preprocessed_df)
                # squarify.plot(sizes=df['SIZE'],
                #               label=df['sql_id'],
                #               color=sb.color_palette("rocket", len(df['SIZE'])), #rocket, magma, Spectral
                #               alpha=.8)
                # plt.axis('off')
                # plt.show()
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._update_tree_map(slide_to_add, preprocessed_df, self.instance_name[i])

    def _update_tree_map(self, slide, df, inst_name):
        """tree_map"""
        bottom_text = self._extract_top_literal_sql_ratio(df, inst_name)

        SlideManager.create_text_box(
            slide,
            self.position["top_menu"]["left_top_text_position"],
            "Top Literal SQL – Top 20 Ratio",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            self.position["top_literal_sql_ratio"]["today_date"],
            self.today_date,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["top_literal_sql_ratio"]["rec_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            slide,
            self.position["top_literal_sql_ratio"]["inst_name_position"],
            inst_name,
            PerformanceAnalyzer.DEFAULT_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        self._set_bottom_text(
            slide,
            bottom_text,
            self.position["top_literal_sql_ratio"]["straight_position"],
            self.position["top_literal_sql_ratio"]["triangle_position"],
            self.position["top_literal_sql_ratio"]["finding_text"],
            self.position["top_literal_sql_ratio"]["bottom_text"],
        )

    def _execute_top_schema_sql(self):
        """
        top_schema_sql에 해당하는 레포트
        """
        self.logger.info("top_schema_sql.pptx")
        top_schema_sql1 = self._convert_sql_to_df(self.sql_path, "TOP_Schema_SQL")
        top_schema_name = top_schema_sql1["parsing_schema_name"].values[0]
        top_schema_sql2 = self._convert_sql_to_df(self.sql_path, "TOP_Schema_SQL2", schema_name=top_schema_name)

        preprocessed_sql2 = self._adjust_df_length(top_schema_sql2, "sql_text")
        self._top_schema_detail(top_schema_sql1, preprocessed_sql2)

    def _top_schema_detail(self, upper_df, down_df):
        """
        top_schema instance_number에 따라 슬라이드 생성
        """
        upper_tp = SlideManager.convert_inches_to_data(self.position["top_schema_sql"]["upper_table"])
        down_tp = SlideManager.convert_inches_to_data(self.position["top_schema_sql"]["down_table"])
        upper_column_width_inches = [Inches(i) for i in self.position["top_schema_sql"]["table1_column_inches"]]
        down_column_width_inches = [Inches(i) for i in self.position["top_schema_sql"]["table2_column_inches"]]

        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.TOP_SCHEMA_SQL)

        for i in range(len(self.instance_number)):
            if i == 0:
                self._detail_top_schema_sql(
                    target_slide,
                    upper_df,
                    down_df,
                    upper_tp,
                    down_tp,
                    upper_column_width_inches,
                    down_column_width_inches,
                    i,
                )

            if i >= 1:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)

                self._detail_top_schema_sql(
                    slide_to_add,
                    upper_df,
                    down_df,
                    upper_tp,
                    down_tp,
                    upper_column_width_inches,
                    down_column_width_inches,
                    i,
                )

    def _arrange_top_schema_sql(self, upper_df, down_df, i):
        up_except_col = ["instance_number", "iowait_time_sec", "iowait_ratio", "rnk"]
        down_except_col = ["instance_number", "rnk"]

        sql1_df = SlideManager.extract_specified_df(
            upper_df, "instance_number", up_except_col, self.instance_number[i]
        ).head(5)

        sql1_df.rename(
            columns={
                "parsing_schema_name": "Schema",
                "elap_time_sec": "Elapsed(sec)",
                "ela_ratio": "Elapsed(%)",
                "cpu_time_sec": "CPU(sec)",
                "cpu_ratio": "CPU(%)",
                "exec_tot": "Exec(count)",
                "exec_ratio": "Exec(%)",
                "lio_tot_block": "Logical I/O(block)",
                "lio_ratio": "Logical(%)",
                "pio_tot_block": "Physical I/O(block)",
                "pio_ratio": "Physical I/O(%)",
            },
            inplace=True,
        )

        sql1_df = sql1_df[
            [
                "Schema",
                "Elapsed(%)",
                "Elapsed(sec)",
                "CPU(%)",
                "CPU(sec)",
                "Exec(%)",
                "Exec(count)",
                "Logical(%)",
                "Logical I/O(block)",
                "Physical I/O(%)",
                "Physical I/O(block)",
            ]
        ]

        sql2_df = SlideManager.extract_specified_df(
            down_df, "instance_number", down_except_col, self.instance_number[i]
        ).head(5)

        sql2_df.rename(
            columns={
                "parsing_schema_name": "Schema",
                "sql_id": "SQL_ID",
                "ela_ratio": "Elapsed(%)",
                "ela_tot": "Elapsed(sec)",
                "cpu_ratio": "CPU(%)",
                "cpu_tot": "CPU(sec)",
                "exec_ratio": "Exec(%)",
                "sql_text": "SQL Text",
            },
            inplace=True,
        )

        sql2_df = sql2_df[
            ["Schema", "SQL_ID", "Elapsed(%)", "Elapsed(sec)", "CPU(%)", "CPU(sec)", "Exec(%)", "SQL Text"]
        ]

        sql1_df = self._set_blank_df(sql1_df, 5)
        sql2_df = self._set_blank_df(sql2_df, 5)

        return sql1_df, sql2_df

    def _detail_top_schema_sql(
        self, slide, upper_df, down_df, upper_tp, down_tp, upper_column_width_inches, down_column_width_inches, i
    ):
        """
        top_schema_sql 공통 부분
        """
        sql1_df, sql2_df = self._arrange_top_schema_sql(upper_df, down_df, i)
        SlideManager.create_text_box(
            slide,
            self.position["top_menu"]["left_top_text_position"],
            "TOP Schema & SQL",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            self.position["top_schema_sql"]["extract_period"],
            self.period,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["top_schema_sql"]["rectangle_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            slide,
            self.position["top_schema_sql"]["oracle_num_position"],
            self.instance_name[i],
            PerformanceAnalyzer.DEFAULT_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        bottom_text = self._extract_top_schema_bottom_text(self.instance_name[i], sql1_df, sql2_df)
        SlideManager.make_table(
            sql1_df,
            upper_tp,
            slide,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            upper_column_width_inches,
        )
        SlideManager.make_table(
            sql2_df,
            down_tp,
            slide,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            down_column_width_inches,
        )

        self._set_bottom_text(
            slide,
            bottom_text,
            self.position["metric"]["straight_position"],
            self.position["metric"]["right_triangle_position"],
            self.position["metric"]["bottom_findings_position"],
            self.position["metric"]["bottom_text_position"],
        )

    def _execute_literal_sql(self):
        """
        literal sql에 해당하는 레포트
        """
        self.logger.info("literalsql.pptx")
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.LITERAL_SQL)
        column_width_inches = [Inches(i) for i in self.position["literal_sql"]["column_width_inches"]]
        tp = SlideManager.convert_inches_to_data(self.position["literal_sql"]["table_position"])

        for i in range(len(self.instance_number)):
            if i == 0:
                self._literal_sql_detail(
                    target_slide, tp, column_width_inches, self.instance_number[i], self.instance_name[i]
                )

            if i >= 1:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._literal_sql_detail(
                    slide_to_add, tp, column_width_inches, self.instance_number[i], self.instance_name[i]
                )

    def _literal_sql_detail(self, slide, tp, column_width_inches, inst_num, instance_name):
        """literal sql 공통 부분"""
        df = self._convert_sql_to_df(self.sql_path, filename="Literal_SQL", inst_num=str(inst_num))
        df = self._arrange_literal_columns(df)
        sp = self.position["bottom_menu"]["bottom_text_position"]

        self._insert_sql_bottom_text(slide, sp, df, instance_name)

        SlideManager.create_text_box(
            slide,
            self.position["top_menu"]["left_top_text_position"],
            "Top Literal SQL (TOP 20 Info)",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            self.position["common_use"]["select_period"],
            f"■ 조회일자 : {self.today_date}",
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["literal_sql"]["ora_rectangle_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            slide,
            self.position["literal_sql"]["ora_position"],
            instance_name,
            PerformanceAnalyzer.DEFAULT_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.make_table(df, tp, slide, Pt(8), Pt(8), column_width_inches)
        self._set_bottom_exclude_text(
            slide,
            self.position["bottom_menu"]["straight_position"],
            self.position["bottom_menu"]["right_triangle_position"],
            self.position["bottom_menu"]["bottom_findings_position"],
        )

    def _arrange_literal_columns(self, df):
        df.rename(
            columns={
                "parsing_schema_name": "SCHEMA",
                "sql_id": "SQL_ID",
                "sql_text": "SQL_TEXT",
                "plan_cnt": "PLAN_CNT",
                "literal_cnt": "LITERAL_CNT",
                "exec_cnt": "EXEC_CNT",
                "elapsed_per": "ELAP_TIME (%)",
                "cpu_time_per": "CPU_TIME (%)",
                "SIZE": "SIZE (%)",
            },
            inplace=True,
        )

        return df

    def _execute_auto_task(self):
        """auto_taske 실행"""
        self.logger.info("auto_task.pptx")
        pd.set_option("display.max_columns", None)
        pd.set_option("display.max_rows", None)
        df1 = self._convert_sql_to_df(self.sql_path, "AUTO_TASK_SQL")
        df2 = self._convert_sql_to_df(self.sql_path, "AUTO_TASK_SQL2")
        df1_val = df1["NVL(AUTOTASK,'(AUTOTOASKSTOP)')"].values[0]

        df2.rename(
            columns={
                "window_name": "WINDOW_NAME",
                "window_start_time": "WINDOW_START_TIME",
                "client_name": "CLENT_NAME",
                "job_start_time": "JOB_START_TIME",
                "job_duration": "JOB_DURATION",
                "job_status": "JOB_STATUS",
            },
            inplace=True,
        )

        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.AUTO_TASK)

        tp2 = SlideManager.convert_inches_to_data(self.position["optimizer_statistics_autotask"]["table_position"])

        SlideManager.create_text_box(
            target_slide,
            self.position["top_menu"]["left_top_text_position"],
            "Optimizer Statistics – AutoTask",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            target_slide,
            self.position["common_use"]["select_period"],
            f"■ 조회일자 : {self.today_date}",
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            target_slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["optimizer_statistics_autotask"]["rectangle_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            target_slide,
            self.position["optimizer_statistics_autotask"]["text1_position"],
            "자동 통계수집 일정",
            Pt(14),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        box_shape = SlideManager.create_shape(
            target_slide,
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            self.position["optimizer_statistics_autotask"]["box_position"],
            PerformanceAnalyzer.LIGHT_GRAY_COLOR,
            PerformanceAnalyzer.LIGHT_GRAY_COLOR,
            0,
        )

        SlideManager.create_text_frame(
            box_shape.text_frame.paragraphs[0],
            df1_val,
            "Roboto Light",
            Pt(11),
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BOLD_FALSE,
            0,
        )

        SlideManager.create_shape(
            target_slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["optimizer_statistics_autotask"]["rectangle2_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            target_slide,
            self.position["optimizer_statistics_autotask"]["text2_position"],
            "자동 통계수집 이력",
            Pt(14),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.make_table(df2, tp2, target_slide, Pt(8), Pt(8))
        SlideManager.optimizer_statistics_style(target_slide)

    def _execute_optimizer_statistics(self):
        """optimizer_statistics실행"""
        self.logger.info("optimizer_statistics.pptx")
        count_df = self._convert_sql_to_df(self.sql_path, "OPTIMIZER_STATISTICS_INFO_SQL1")
        avg_df = self._convert_sql_to_df(self.sql_path, "OPTIMIZER_STATISTICS_INFO_SQL2")

        left_chart_position = self.position["optimizer_statistics"]["left_chart_position"]
        right_chart_position = self.position["optimizer_statistics"]["right_chart_position"]

        slide_num, target_slide = SlideManager.read_slide(
            self.presentation.slides, PerformanceAnalyzer.OPTIMIZER_STATISTICS
        )

        SlideManager.create_text_box(
            target_slide,
            self.position["top_menu"]["left_top_text_position"],
            "Optimizer Statistics",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        count_df["retention"] = ["less month", "less year", "over year", "null", "0"]
        count_df_category = count_df["retention"].tolist()
        count_df_val = tuple(count_df["cnt"].tolist())

        avg_df_category = avg_df["owner"].tolist()
        avg_df_val = tuple(avg_df["avg_day"].tolist())

        SlideManager.add_bar_chart(count_df_category, count_df_val, left_chart_position, target_slide)
        SlideManager.add_bar_chart(avg_df_category, avg_df_val, right_chart_position, target_slide)

        SlideManager.set_bar_chart_style(target_slide)

        SlideManager.create_shape(
            target_slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["optimizer_statistics"]["rectangle_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            target_slide,
            self.position["optimizer_statistics"]["text1_position"],
            "오브젝트 현황",
            Pt(13),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            target_slide,
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            self.position["optimizer_statistics"]["rectangle2_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            target_slide,
            self.position["optimizer_statistics"]["text2_position"],
            "스키마 현황",
            Pt(13),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            target_slide,
            self.position["common_use"]["select_period"],
            f"■ 조회일자 : {self.today_date}",
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        bottom_text = self._extract_optimizer_statistics_bottom_text(count_df, avg_df)
        self._set_bottom_text(
            target_slide,
            bottom_text,
            self.position["optimizer_statistics"]["straight_position"],
            self.position["optimizer_statistics"]["right_triangle_position"],
            self.position["optimizer_statistics"]["bottom_findings_position"],
            self.position["optimizer_statistics"]["bottom_text_position"],
        )

    def _execute_optimizer_statistics_zero_null(self):
        self.logger.info("optimizer_statistics_zero_null.pptx")

        left_tp = SlideManager.convert_inches_to_data(self.position["optimizer_statistics_zero_null"]["left_table"])
        right_tp = SlideManager.convert_inches_to_data(self.position["optimizer_statistics_zero_null"]["right_table"])
        column_width_inches = [
            Inches(i) for i in self.position["optimizer_statistics_zero_null"]["column_width_inches"]
        ]

        df1 = self._convert_sql_to_df(self.sql_path, "OPTIMIZER_STATISTICS_ZERO_NULL_SQL1")
        df1 = self._set_blank_df(df1, 20)

        df2 = self._convert_sql_to_df(self.sql_path, "OPTIMIZER_STATISTICS_ZERO_NULL_SQL2")
        df2 = self._set_blank_df(df2, 20)

        df1.columns = map(lambda x: x.upper(), df1.columns)
        df2.columns = map(lambda x: x.upper(), df2.columns)

        slide_num, target_slide = SlideManager.read_slide(
            self.presentation.slides, PerformanceAnalyzer.OPTIMIZER_STATISTICS_ZERO_NULL
        )
        SlideManager.create_text_box(
            target_slide,
            self.position["top_menu"]["left_top_text_position"],
            "Optimizer Statistics – Statistics null or 0",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            target_slide,
            self.position["optimizer_statistics_zero_null"]["today_date"],
            f"■ 조회일자 : {self.today_date}",
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            target_slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["optimizer_statistics_zero_null"]["rectangle_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            target_slide,
            self.position["optimizer_statistics_zero_null"]["text1_position"],
            "통계정보 NULL 테이블 리스트",
            Pt(11),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            target_slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["optimizer_statistics_zero_null"]["rectangle2_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            target_slide,
            self.position["optimizer_statistics_zero_null"]["text2_position"],
            "통계정보 NUM_ROWS=0  테이블 리스트",
            Pt(11),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.make_table(df1, left_tp, target_slide, Pt(10), Pt(9), column_width_inches)
        SlideManager.make_table(df2, right_tp, target_slide, Pt(10), Pt(9), column_width_inches)

    def _top3_chart_detail(self, slide, instance_num, sql1_df_mspsec):
        """top_3_wait_event에 있는 차트 스타일"""
        unit_tuple = SlideManager.convert_unit(tuple(sql1_df_mspsec.tolist()))
        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                SlideManager.set_chart_style(chart)
                plot = chart.plots[0]

                for indx, series in enumerate(plot.series):
                    for point_indx, label in enumerate(series.points):
                        color_indx = int(instance_num) - 1
                        color = PerformanceAnalyzer.COLOR_LIST[color_indx]

                        value_to_set = unit_tuple[point_indx]
                        tf = label.data_label.text_frame
                        tf.text = value_to_set

                        for paragraph in tf.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(8)
                                run.font.color.rgb = PerformanceAnalyzer.COLOR_LIST[color_indx]
                                run.font.name = PerformanceAnalyzer.FONT_NAME
                                run.font.bold = PerformanceAnalyzer.BOLD_TRUE

                        SlideManager.set_label_style(
                            label,
                            color,
                            Pt(8),
                            PerformanceAnalyzer.FONT_NAME,
                            True,
                            XL_DATA_LABEL_POSITION.ABOVE,
                        )
                        SlideManager.set_chart_marker_style(
                            series,
                            XL_MARKER_STYLE.CIRCLE,
                            color,
                            PerformanceAnalyzer.WHILE_COLOR,
                        )
                        SlideManager.set_chart_line_style(series, color, Pt(2.25))

    def _extract_top_seo_io_df(self):
        if len(self.instance_number) == 1:
            df = self._convert_sql_to_df(self.sql_path, "TOP_SEG_IO(SINGLE)")
            df = df.rename(
                columns={
                    "name": "NAME",
                    "logical_reads_ratio": "LOGICAL_READS(%)",
                    "physical_reads_ratio": "PHYSICAL_READS(%)",
                    "physical_writes_ratio": "PHYSICAL_WRITES(%)",
                    "row_lock_waits_ratio": "ROW_LOCK_WAITS(%)",
                    "buffer_busy_waits_ratio": "BUFFER_BUSY_WAIT(%)",
                }
            )

        else:
            df = self._convert_sql_to_df(self.sql_path, "TOP_SEG_IO(RAC)")
            df = df.rename(
                columns={
                    "name": "NAME",
                    "logical_reads_ratio": "LOGICAL_READS(%)",
                    "physical_reads_ratio": "PHYSICAL_READS(%)",
                    "physical_writes_ratio": "PHYSICAL_WRITES(%)",
                    "row_lock_waits_ratio": "ROW_LOCK_WAITS(%)",
                    "gc_buffer_busy_ratio": "GC_BUFFER_BUSY(%)",
                }
            )

        return df

    def _execute_top_seg_io_block(self):
        self.logger.info("top_seg_io_block.pptx")

        left_chart_position = self.position["top_seg_block"]["left_chart_position"]
        right_chart_position = self.position["top_seg_block"]["right_chart_position"]

        df = self._convert_sql_to_df(self.sql_path, "TOP_SEG_IO(BLOCK)")
        slide_num, target_slide = SlideManager.read_slide(
            self.presentation.slides, PerformanceAnalyzer.TOP_SEGMENTS_BLOCK
        )

        for i in range(len(self.instance_name)):
            if i == 0:
                self._top_seg_io_block_detail(
                    df,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    target_slide,
                    left_chart_position,
                    right_chart_position,
                )

            elif i > 1 and i % 2 == 0:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._top_seg_io_block_detail(
                    df,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    slide_to_add,
                    left_chart_position,
                    right_chart_position,
                )

    def _top_seg_io_block_detail(self, df, inst_number_range, inst_name_range, slide, left_pt, right_pt):
        SlideManager.create_text_box(
            slide,
            self.position["top_menu"]["left_top_text_position"],
            "Top Segments I/O – Block",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            self.position["common_use"]["extract_period"],
            self.period,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        bottom_text = self._extract_top_segments_block_bottom(df, inst_number_range, inst_name_range)
        self._set_bottom_text(
            slide,
            bottom_text,
            self.position["bottom_menu"]["straight_position"],
            self.position["bottom_menu"]["right_triangle_position"],
            self.position["bottom_menu"]["bottom_findings_position"],
            self.position["bottom_menu"]["bottom_text_position"],
        )

        for i in range(len(inst_number_range)):
            rec_position = (
                self.position["top_seg_block"]["left_rectangle_position"]
                if self.is_even(i)
                else self.position["top_seg_block"]["right_rectangle_position"]
            )

            instance_name_position = (
                self.position["top_seg_block"]["text1_position"]
                if self.is_even(i)
                else self.position["top_seg_block"]["text2_position"]
            )

            pt = left_pt if self.is_even(i) else right_pt

            inst_df = df[df["instance_number"] == inst_number_range[i]]
            inst_df = inst_df.drop(columns=["instance_number", "physical_writes_block"])

            category_df = inst_df["name"].tolist()
            df_val1 = inst_df["logical_reads_block"]
            df_val2 = inst_df["physical_reads_block"]
            max_value = inst_df["logical_reads_block"].max()

            SlideManager.create_shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                rec_position,
                PerformanceAnalyzer.YELLOW_COLOR,
                PerformanceAnalyzer.YELLOW_COLOR,
                0,
            )

            SlideManager.create_text_box(
                slide,
                instance_name_position,
                inst_name_range[i],
                Pt(14),
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_TRUE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            chart = SlideManager.add_stacked_two_bar_chart(category_df, df_val1, df_val2, pt, slide)
            y_axis_max_value = SlideManager.make_max_value([max_value])
            label_num_format = SlideManager.set_y_axis_max_value(chart, y_axis_max_value)
            SlideManager.set_stacked_bar_style(slide, label_num_format)

    def _execute_top_seg_io_ratio(self):
        self.logger.info("top_seg_io_ratio.pptx")

        left_chart_position = self.position["top_seg_block"]["left_chart_position"]
        right_chart_position = self.position["top_seg_block"]["right_chart_position"]

        df = self._convert_sql_to_df(self.sql_path, "TOP_SEG_IO(RATIO)")
        slide_num, target_slide = SlideManager.read_slide(
            self.presentation.slides, PerformanceAnalyzer.TOP_SEGMENTS_RATIO
        )

        for i in range(len(self.instance_name)):
            if i == 0:
                self._top_seg_io_ratio_detail(
                    df,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    target_slide,
                    left_chart_position,
                    right_chart_position,
                )

            elif i > 1 and i % 2 == 0:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._top_seg_io_ratio_detail(
                    df,
                    self.instance_number[i : i + 2],
                    self.instance_name[i : i + 2],
                    slide_to_add,
                    left_chart_position,
                    right_chart_position,
                )

    def _top_seg_io_ratio_detail(self, df, inst_number_range, inst_name_range, slide, left_pt, right_pt):
        SlideManager.create_text_box(
            slide,
            self.position["top_menu"]["left_top_text_position"],
            "Top Segments I/O – Ratio",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            self.position["common_use"]["extract_period"],
            self.period,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        bottom_text = self._extract_top_segments_ratio_bottom(df, inst_number_range, inst_name_range)
        self._set_bottom_text(
            slide,
            bottom_text,
            self.position["bottom_menu"]["straight_position"],
            self.position["bottom_menu"]["right_triangle_position"],
            self.position["bottom_menu"]["bottom_findings_position"],
            self.position["bottom_menu"]["bottom_text_position"],
        )

        for i in range(len(inst_number_range)):
            rec_position = (
                self.position["top_seg_block"]["left_rectangle_position"]
                if self.is_even(i)
                else self.position["top_seg_block"]["right_rectangle_position"]
            )

            instance_name_position = (
                self.position["top_seg_block"]["text1_position"]
                if self.is_even(i)
                else self.position["top_seg_block"]["text2_position"]
            )

            pt = left_pt if self.is_even(i) else right_pt

            inst_df = df[df["instance_number"] == inst_number_range[i]]
            inst_df = inst_df.drop(columns=["instance_number"])

            preprocessed_df = inst_df.copy()
            preprocessed_df["logical_reads_ratio"] = (
                preprocessed_df["logical_reads_ratio"].replace("%", "", regex=True).astype(float)
            )
            preprocessed_df["physical_reads_ratio"] = (
                preprocessed_df["physical_reads_ratio"].replace("%", "", regex=True).astype(float)
            )
            preprocessed_df["physical_writes_ratio"] = (
                preprocessed_df["physical_writes_ratio"].replace("%", "", regex=True).astype(float)
            )

            category_df = preprocessed_df["name"].tolist()
            df_val1 = preprocessed_df["logical_reads_ratio"]
            df_val2 = preprocessed_df["physical_reads_ratio"]
            df_val3 = preprocessed_df["physical_writes_ratio"]

            SlideManager.create_shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                rec_position,
                PerformanceAnalyzer.YELLOW_COLOR,
                PerformanceAnalyzer.YELLOW_COLOR,
                0,
            )

            SlideManager.create_text_box(
                slide,
                instance_name_position,
                inst_name_range[i],
                Pt(14),
                PerformanceAnalyzer.FONT_NAME,
                PerformanceAnalyzer.BOLD_TRUE,
                PerformanceAnalyzer.BLACK_COLOR,
                PerformanceAnalyzer.BASE_LINE_SPACE,
            )

            SlideManager.add_stacked_three_bar_chart(category_df, df_val1, df_val2, df_val3, pt, slide)
            SlideManager.set_stacked_bar_style(slide, '0"%"')

    def _execute_top_seg_io_table(self):
        self.logger.info("top_seg_io_table.pptx")
        df = self._extract_top_seo_io_df()

        df = df.applymap(lambda x: float(x.replace("%", "")) if isinstance(x, str) and "%" in str(x) else x)
        df = df.applymap(lambda x: f"{x:.1f} %" if x != 0.0 and isinstance(x, float) else x)
        df = df.applymap(lambda x: "0 %" if x == 0.0 and isinstance(x, float) else x)

        tp = SlideManager.convert_inches_to_data(self.position["top_seg_io"]["table_position"])
        column_width_inches = [Inches(i) for i in self.position["top_seg_io"]["column_width_inches"]]
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.TOP_SEG_IO)

        for indx, inst_num in enumerate(self.instance_number):
            inst_df = df[df["instance_number"] == inst_num]
            inst_num_drop_df = inst_df.drop(columns=["instance_number"])

            if indx == 0:
                self._top_seg_io_detail(
                    target_slide, inst_num_drop_df, self.instance_name[indx], tp, column_width_inches
                )

            else:
                slide_num, slide_to_add = SlideManager.add_slide(self.presentation, slide_num)
                self._top_seg_io_detail(
                    slide_to_add, inst_num_drop_df, self.instance_name[indx], tp, column_width_inches
                )

    def _top_seg_io_detail(self, slide, df, inst_name, tp, column_width_inches):
        bottom_text = self._extract_top_seg_io_bottom_text(df, inst_name)

        SlideManager.create_text_box(
            slide,
            self.position["top_menu"]["left_top_text_position"],
            "Top Segments I/O - Detail",
            PerformanceAnalyzer.LARGE_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            self.position["common_use"]["extract_period"],
            self.period,
            PerformanceAnalyzer.SMALL_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.position["top_seg_io"]["rectangle_position"],
            PerformanceAnalyzer.YELLOW_COLOR,
            PerformanceAnalyzer.YELLOW_COLOR,
            0,
        )

        SlideManager.create_text_box(
            slide,
            self.position["top_seg_io"]["inst_name_position"],
            inst_name,
            Pt(14),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.make_table(df, tp, slide, Pt(10), Pt(9), column_width_inches)
        self._set_bottom_text(
            slide,
            bottom_text,
            self.position["bottom_menu"]["straight_position"],
            self.position["bottom_menu"]["right_triangle_position"],
            self.position["bottom_menu"]["bottom_findings_position"],
            self.position["bottom_menu"]["bottom_text_position"],
        )

    def _extract_top_seg_io_bottom_text(self, df, inst_name):
        """top_seg_io_bottom_text 추출"""
        logical_val = df["LOGICAL_READS(%)"].values[0]
        object_name = df["NAME"].values[0]

        text = f"{inst_name}에서  {object_name} 세그먼트가 인스턴스의 Logical Reads {logical_val}를 차지함. 집중적인 I/O발생에 대한 튜닝 필요"

        return text

    def _extract_top3_list(self, df):
        """
        top_n에서 인스턴스별로 top3 추출
        """
        instance_dict = {}

        for instance_num in df["instance_number"].unique():
            instance_num_df = df[df["instance_number"] == instance_num]
            top_3_df = instance_num_df[instance_num_df["event_name"] != "DB CPU"].head(3)
            top3_3_name = list(top_3_df["event_name"])
            instance_num_name = f"instance_{instance_num}"
            instance_dict[instance_num_name] = top3_3_name

        return instance_dict

    def _extract_preprocessed_df(self, event_name, instance_num, sql_filename, except_col=None):
        """top3에 해당하는 df 전처리"""
        if except_col is None:
            except_col = []

        df = self._convert_sql_to_df(self.sql_path, sql_filename, event_name)
        preprocessed_df = self._set_df_date_time(df)
        event_name_df = preprocessed_df[preprocessed_df["instance_number"] == int(instance_num)]
        event_name_df = event_name_df.drop(columns=except_col)
        return event_name_df

    def _set_bottom_exclude_text(self, slide, straight_position, triangle_position, finding_position):
        """슬라이드 하단에 공통으로 기입된 텍스트 및 도형"""
        SlideManager.create_connector(slide, straight_position)

        SlideManager.create_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
            triangle_position,
            PerformanceAnalyzer.DARK_GRAY_COLOR,
            PerformanceAnalyzer.DARK_GRAY_COLOR,
            -90,
        )

        SlideManager.create_text_box(
            slide,
            finding_position,
            "Findings",
            PerformanceAnalyzer.DEFAULT_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

    def _set_bottom_text(
        self, slide, bottom_text, straight_position, triangle_position, finding_position, text_position
    ):
        """슬라이드 하단에 공통으로 기입된 텍스트 및 도형"""
        SlideManager.create_connector(slide, straight_position)

        SlideManager.create_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
            triangle_position,
            PerformanceAnalyzer.DARK_GRAY_COLOR,
            PerformanceAnalyzer.DARK_GRAY_COLOR,
            -90,
        )

        SlideManager.create_text_box(
            slide,
            finding_position,
            "Findings",
            PerformanceAnalyzer.DEFAULT_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.BASE_LINE_SPACE,
        )

        SlideManager.create_text_box(
            slide,
            text_position,
            bottom_text,
            PerformanceAnalyzer.DEFAULT_FONT_SIZE,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_FALSE,
            PerformanceAnalyzer.BLACK_COLOR,
            PerformanceAnalyzer.LONG_LINE_SPACE,
        )

    def _chart_style(self, slide, up_col_tuple, down_col_tuple):
        """chart_style"""
        chart_list = []

        for indx, shape in enumerate(slide.shapes):
            if shape.has_chart:
                chart = shape.chart
                SlideManager.set_chart_style(chart)
                chart_list.append(chart)

        self._set_chart_label(chart_list, up_col_tuple, down_col_tuple)

    def _set_chart_label(self, chart_list, up_col_tuple, down_col_tuple):
        upper_plot = chart_list[0].plots[0]
        down_plot = chart_list[1].plots[0]

        SlideManager.create_label_text(
            upper_plot,
            up_col_tuple,
            PerformanceAnalyzer.COLOR_LIST,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
        )

        SlideManager.create_label_text(
            down_plot,
            down_col_tuple,
            PerformanceAnalyzer.COLOR_LIST,
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
        )

    def _chart_style_detail(self, series, label, color):
        """chart_style: label, line, marker"""
        SlideManager.set_label_style(
            label,
            color,
            Pt(8),
            PerformanceAnalyzer.FONT_NAME,
            PerformanceAnalyzer.BOLD_TRUE,
            XL_DATA_LABEL_POSITION.ABOVE,
        )
        SlideManager.set_chart_line_style(series, color, Pt(2.25))
        SlideManager.set_chart_marker_style(series, XL_MARKER_STYLE.CIRCLE, color, PerformanceAnalyzer.WHILE_COLOR)

    def _set_df_date_time(self, df):
        """DATE_TIME 컬럼에 해당하는 데이터 수정"""
        if "mmdd" in df.columns and "hh24" in df.columns:
            df["hh24"] = df["hh24"].astype(str) + ":00"
            df.insert(loc=0, column="combined_datetime", value="")
            df["combined_datetime"] = df["mmdd"] + " " + df["hh24"]
            df = df.drop(["mmdd", "hh24"], axis=1)
            df["combined_datetime"] = df["combined_datetime"].apply(lambda x: "\n".join(x.split()))

        return df

    def _adjust_df_length(self, df, col):
        if not (df[col] == "").any():
            df[col] = df[col].str.replace(r"\s+", " ", regex=True).apply(lambda x: x.lstrip())
            df[col] = df[col].str.slice(0, 55)
            df[col] = df[col].astype(str) + "..."

        return df

    def _convert_max_unit(self, score):
        if 4 <= len(str(score)) < 7:
            new_score = f"{score/1000:.1f}K"

        elif len(str(score)) >= 7:
            new_score = f"{score / 1000000:.1f}M"

        else:
            new_score = score

        return new_score

    def _extract_chart_bottom_text(self, slide, metric_name, df, col1, col2, val1, val2):
        bottom_text = ""

        for idx, instance_num in enumerate(self.instance_number):
            extract_df = df[df["instance_number"] == instance_num]
            inst_name = self.instance_name[idx]

            avg_score = extract_df[col1].max()
            avg_new_score = self._convert_max_unit(avg_score)
            max_score = extract_df[col2].max()
            max_new_score = self._convert_max_unit(max_score)

            instance_text = f"{inst_name}의 {metric_name} {val1} 최고값은 {avg_new_score}이며, {metric_name} {val2}의 최고값은 {max_new_score}로 확인.\n"
            bottom_text += instance_text

        return bottom_text

    def _extract_top_n_bottom_text(self, top3_inst_dict, instance_num_range, instance_name_range):
        """top_n_bottom_text 추출"""
        bottom_text = ""
        for idx, inst_num in enumerate(instance_num_range):
            top3_inst_dict_key = f"instance_{inst_num}"
            top3_inst_dict_val = top3_inst_dict[top3_inst_dict_key]
            top1_val = top3_inst_dict_val[0]
            top2_val = top3_inst_dict_val[1]
            top3_val = top3_inst_dict_val[2]

            inst_name = instance_name_range[idx]
            instance_text = f"{inst_name}에서 {top1_val}, {top2_val}, {top3_val} 이벤트가 가장 많이 발생함.\n"

            bottom_text += instance_text
        return bottom_text

    def _extract_y_axis_value(self, metric_name):
        if metric_name not in PerformanceAnalyzer.y_axis_dict.keys():
            PerformanceAnalyzer.y_axis_dict[metric_name] = "Count (K=1000)"

        return PerformanceAnalyzer.y_axis_dict[metric_name]

    def _extract_top3_text(self, inst_name, event_name, sql1_df, sql2_df):
        max_mspsec_value = sql1_df["mspsec"].max()
        max_date_time = sql1_df.loc[sql1_df["mspsec"] == max_mspsec_value, "combined_datetime"].values[0]
        split_date = max_date_time.split("\n")[0].split("-")
        joined_result = "".join(split_date)
        date = f"{self.config['args']['s_date'][:4]}{joined_result}"
        date = datetime.strptime(date, "%Y%m%d")
        date = date.strftime("%Y.%m.%d")

        avg_wait = sql2_df["avg_wait"].values[0]

        split_time = int(max_date_time.split("\n")[1].split(":")[0])
        before_time = str(split_time - 1).zfill(2)
        after_time = str(split_time + 1).zfill(2)

        bottom_text = (
            f"{inst_name}에서 {event_name}이벤트가 가장 많이 발생함.\n"
            f"{date},{before_time}~{after_time}시에 {avg_wait} 발생하였으며, 해당시간에 집중적으로 발생함."
        )

        return bottom_text

    def _extract_top_schema_bottom_text(self, inst_name, sql1, sql2):
        elapsed_max_value = sql1.loc[0, "Elapsed(%)"]
        sql1_schema_name = sql1.loc[0, "Schema"]
        sql2_schema_name = sql2.loc[0, "Schema"]

        bottom_text = (
            f"{inst_name}에서 {sql1_schema_name} 스키마의 수행시간이 전체 {elapsed_max_value}%를 차지.\n"
            f"{sql2_schema_name} 스키마에서 가장 많이 수행된 SQL은 위와 같음."
        )

        return bottom_text

    def _set_blank_df(self, df, num):
        df_append = pd.DataFrame("", index=range(num - len(df)), columns=df.columns)
        df = pd.concat([df, df_append], ignore_index=True)
        return df

    def _insert_memory_bottom_text(self, slide, sp, green_text, common_text):
        tb = slide.shapes.add_textbox(sp["left"], sp["top"], sp["width"], sp["height"])
        tf = tb.text_frame

        paragraph = tf.paragraphs[0]
        paragraph.font.size = Pt(10)
        paragraph.line_spacing = 1.5

        run = paragraph.add_run()
        run.text = green_text
        run.font.color.rgb = PerformanceAnalyzer.LIGHT_GREEN_COLOR

        paragraph_2 = tf.add_paragraph()
        paragraph_2.text = common_text
        paragraph_2.font.size = Pt(10)
        paragraph_2.line_spacing = 1.5

    def _insert_sql_bottom_text(self, slide, sp, df, inst_name):
        size_val = df["SIZE (%)"].str.replace("%", "").astype(float)
        size_sum = size_val.sum()

        bottom_text1 = f"{inst_name}는 Top 20 Literal SQL이 전체 공간 중 "
        bottom_text2 = f"{size_sum:.1f}%"
        bottom_text3 = "비율을 차지함. 바인드 변수 사용 권고."

        tb = slide.shapes.add_textbox(sp["left"], sp["top"], sp["width"], sp["height"])
        tf = tb.text_frame

        paragraph = tf.paragraphs[0]
        paragraph.font.size = Pt(10)

        run = paragraph.add_run()
        run.text = bottom_text1

        run2 = paragraph.add_run()
        run2.text = bottom_text2
        font2 = run2.font
        font2.color.rgb = RGBColor(255, 0, 0)

        run3 = paragraph.add_run()
        run3.text = bottom_text3

    def _extract_optimizer_statistics_bottom_text(self, count_df, avg_df):
        less_month_cnt = count_df[count_df["retention"] == "less month"]["cnt"].values[0]
        less_year_cnt = count_df[count_df["retention"] == "less year"]["cnt"].values[0]
        one_year_cnt = count_df[count_df["retention"] == "over year"]["cnt"].values[0]
        null_cnt = count_df[count_df["retention"] == "null"]["cnt"].values[0]
        zero_cnt = count_df[count_df["retention"] == "0"]["cnt"].values[0]
        total_avg = avg_df["avg_day"].mean()

        bottom_text = (
            f"ORCLDB의 오브젝트별 통계정보 기간은 1달 이내 {less_month_cnt}개, 1년 이내 {less_year_cnt}개, "
            f"1년 이상 {one_year_cnt}개, 미 수집은{null_cnt}개, 0으로 수집은 {zero_cnt}개로 확인됨.\n"
            f"ORCLDB의 평균 통계수집일은 최대 {total_avg}일로 확인됨\n\n"
            f"※ 통계정보는 Database의 SQL 성능과 관련이 깊음.\n"
            f"※ 정확한 통계정보일수록 Optimizer는 최적의 실행계획을 수립할 수 있음."
        )

        return bottom_text

    def _extract_top_segments_block_bottom(self, df, inst_num_range, inst_name_range):
        bottom_text = ""

        for idx, inst_num in enumerate(inst_num_range):
            inst_df = df[df["instance_number"] == inst_num]
            top_segment = inst_df["logical_reads_block"].max()
            index_of_max_value = inst_df[inst_df["logical_reads_block"] == top_segment].index[0]
            top_name = inst_df.loc[index_of_max_value, "name"]
            inst_name = inst_name_range[idx]
            instance_text = f"{inst_name}에서 {top_name} 세그먼트가 가장 많은 Logical reads 수행. 집중적인 I/O발생에 대한 튜닝 필요\n"
            bottom_text += instance_text

        return bottom_text

    def _extract_top_segments_ratio_bottom(self, df, inst_num_range, inst_name_range):
        bottom_text = ""

        for idx, inst_num in enumerate(inst_num_range):
            inst_df = df[df["instance_number"] == inst_num]
            top_segment = inst_df["logical_reads_ratio"].max()
            index_of_max_value = inst_df[inst_df["logical_reads_ratio"] == top_segment].index[0]
            top_name = inst_df.loc[index_of_max_value, "name"]
            inst_name = inst_name_range[idx]
            instance_text = (
                f"{inst_name}에서 {top_name} 세그먼트가 인스턴스의 Logical reads {top_segment}를 차지함. 집중적인 I/O 발생에 대한 튜닝 필요\n"
            )
            bottom_text += instance_text

        return bottom_text

    def _extract_memory_bottom_text(self, inst_num_range, inst_name_range, text1, text2, df, col1, col2):
        bottom_text = ""

        sga_size_list = df[col1].tolist()
        physical_reads_list = df[col2].tolist()

        for idx, inst_num in enumerate(inst_num_range):
            sga_size = sga_size_list[idx]
            physical_reads = physical_reads_list[idx]
            new_physical_reads = (
                f"0{physical_reads}" if "." in physical_reads and physical_reads[0] == "." else physical_reads
            )
            inst_name = inst_name_range[idx]

            instance_text = (
                f"{inst_name}의 {text1} 영역을 {sga_size}로 증가할 경우 {text2}가 현재 대비 {new_physical_reads}의 개선이 될 것으로 예상됨\n"
            )
            bottom_text += instance_text

        return bottom_text

    def _extract_top_literal_sql_ratio(self, df, inst_name):
        sum_of_values = df.loc[1:, "SIZE"].sum()
        sum_df_values_rounded = round(sum_of_values, 1)
        bottom_text = f"{inst_name}는 Top 20 Literal SQL이 전체 공간 중 {sum_df_values_rounded}% 비율을 차지함. 바인드 SQL 사용 권고."

        return bottom_text
