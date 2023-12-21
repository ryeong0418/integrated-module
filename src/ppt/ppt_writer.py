from src.ppt.pd2ppt import df_to_table
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.enum.chart import XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.chart import XL_MARKER_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement
import re


class SlideManager:
    """
    SlideManager class
    """

    @staticmethod
    def convert_inches_to_data(dict_position):
        """convert inches data"""
        for k, v in dict_position.items():
            dict_position[k] = Inches(v)
        return dict_position

    @staticmethod
    def read_slide(slides, text_frame_text):
        """read specific slide"""
        for idx, slide in enumerate(slides):
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text == text_frame_text:
                    sp = shape._element  # element
                    sp.getparent().remove(sp)
                    return (idx, slide)

    @staticmethod
    def delete_slide(prs):
        """delete specific slide"""
        xml_slides = prs.slides._sldIdLst
        xml_slides_list = list(xml_slides)
        for idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame and re.search(r"^#", shape.text_frame.text):
                    xml_slides.remove(xml_slides_list[idx])

    @staticmethod
    def add_slide(prs, idx):
        """add slide needed"""
        source_slide = prs.slides[idx]
        slide_layout = source_slide.slide_layout
        copied_slide = prs.slides.add_slide(slide_layout)
        xml_slides = prs.slides._sldIdLst
        idx += 1

        for shape in copied_slide.placeholders:
            title_placeholder = copied_slide.placeholders[shape.placeholder_format.idx]
            sp_title = title_placeholder.element
            sp_title.getparent().remove(sp_title)
        xml_slides.insert(idx, xml_slides[-1])

        return idx, copied_slide

    @staticmethod
    def extract_specified_df(df, col_name, except_col, i):
        """특정 dataframe 추출"""
        df = df[df[col_name] == i]
        df = df.drop(columns=except_col)
        return df

    @staticmethod
    def create_text_box(slide, sp, content_text, font_size, font_name, bold_effect, font_color, line_space):
        """create textbox, text"""
        tb = slide.shapes.add_textbox(sp["left"], sp["top"], sp["width"], sp["height"])
        tf = tb.text_frame

        paragraph = tf.paragraphs[0]
        paragraph.text = content_text
        paragraph.font.size = font_size
        paragraph.font.bold = bold_effect
        paragraph.font.name = font_name
        paragraph.font.color.rgb = font_color
        paragraph.line_spacing = line_space

    # @staticmethod
    # def create_two_color_text_box(slide, sp, content_text, font_size, font_name, bold_effect, line_space):
    #     tb = slide.shapes.add_textbox(sp["left"], sp["top"], sp["width"], sp["height"])
    #     tf = tb.text_frame
    #
    #     paragraph = tf.paragraphs[0]
    #     run = paragraph.add_run()
    #
    #     run.characters[0:5].font.color.rgb = RGBColor(255, 0, 0)  # 빨간색(RGB)
    #
    #     # 나머지 텍스트에 대한 서식을 지정할 수 있습니다.
    #     run.font.size = font_size
    #     run.font.bold = bold_effect
    #     run.font.name = font_name
    #     run.line_spacing = line_space
    #
    #     # paragraph.text = content_text
    #     # paragraph.font.size = font_size
    #     # paragraph.font.bold = bold_effect
    #     # paragraph.font.name = font_name
    #     # paragraph.line_spacing = line_space

    @staticmethod
    def create_text_frame(paragraph, content_text, font_name, font_size, font_color, font_bold, line_space):
        """create text"""
        paragraph.text = content_text
        paragraph.font.name = font_name
        paragraph.font.size = font_size
        paragraph.font.bold = font_bold
        paragraph.font.color.rgb = font_color
        paragraph.line_spacing = line_space

    @staticmethod
    def create_label_text(plot, col_tuple, color_list, font_name, bold):
        """label을 원하는 형태로 출력하기 위해 필요한 함수"""
        for index, series in enumerate(plot.series):
            for point_indx, label in enumerate(series.points):
                value_to_set = col_tuple[index][point_indx]
                tf = label.data_label.text_frame
                tf.text = value_to_set

                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(8)
                        run.font.color.rgb = color_list[index]
                        run.font.name = font_name
                        run.font.bold = bold

                label.data_label.position = XL_DATA_LABEL_POSITION.ABOVE
                SlideManager.set_chart_line_style(series, color_list[index], Pt(2.25))
                SlideManager.set_chart_marker_style(
                    series, XL_MARKER_STYLE.CIRCLE, color_list[index], RGBColor(255, 255, 255)
                )

    @staticmethod
    def create_shape(slide, type_shape, sp, fill_color, line_color, rt):
        """create shape"""
        shapes = slide.shapes
        shape = shapes.add_shape(type_shape, sp["left"], sp["top"], sp["width"], sp["height"])

        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = fill_color

        shape_line = shape.line.fill
        shape_line.solid()
        shape_line.fore_color.rgb = line_color
        shape.rotation = rt

        return shape

    @staticmethod
    def create_connector(slide, sp):
        """create connector"""
        line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, sp["left"], sp["top"], sp["width"], sp["height"])
        line.rotation = 0
        line.line.dash_style = MSO_LINE_DASH_STYLE.SQUARE_DOT
        line.line.color.rgb = RGBColor(157, 157, 157)
        line.line.width = Pt(0.75)

    @staticmethod
    def make_max_value(max_list):
        """
        y축의 최대값을 만든다.
        예를 들어, 출력된 최대값이 385일때, 최대값을 400으로 반환하고,
                 출력된 최대값이 1101일때, 최대값을 2000으로 반환한다.
        """
        if not max_list:
            return 10

        max_value = int(max(max_list))

        if 0 <= max_value <= 9:
            max_value = 10

        else:
            num_place = 10 ** (len(str(max_value)) - 1)
            max_value = ((max_value // num_place) + 1) * num_place

        return int(max_value)

    @staticmethod
    def insert_chart(slide, df_category, df_val, series_name, pt):
        """라인 차트 insert"""
        chart_data = CategoryChartData()
        chart_data.categories = df_category
        chart_data.add_series(series_name, df_val)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, pt["x"], pt["y"], pt["cx"], pt["cy"], chart_data).chart

        return chart

    @staticmethod
    def insert_line_chart(slide, df, pt, inst_num_range, inst_name_range, col):
        """라인 차트 insert"""
        chart_data = CategoryChartData()

        for idx, inst_num in enumerate(inst_num_range):
            extract_df = df[df["instance_number"] == inst_num]
            chart_categories = extract_df["combined_datetime"].tolist()

            col_value_tuple = tuple(extract_df[col].tolist())
            chart_data.categories = chart_categories
            chart_data.add_series(inst_name_range[idx], col_value_tuple)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, pt["x"], pt["y"], pt["cx"], pt["cy"], chart_data).chart

        return chart

    @staticmethod
    def add_bar_chart(df_category, df_val, tp, slide, series1=""):
        """한 개 막대그래프 추가"""
        chart_data = CategoryChartData()
        chart_data.categories = df_category
        chart_data.add_series(series1, df_val)

        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, tp["x"], tp["y"], tp["cx"], tp["cy"], chart_data)

    @staticmethod
    def add_close_bar_chart(df_category, df_val1, df_val2, series1, series2, tp, slide):
        """두 개 막대그래프 추가"""
        chart_data = CategoryChartData()
        chart_data.categories = df_category
        chart_data.add_series(series1, df_val1)
        chart_data.add_series(series2, df_val2)

        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, tp["x"], tp["y"], tp["cx"], tp["cy"], chart_data)

    @staticmethod
    def add_stacked_two_bar_chart(df_category, df_val1, df_val2, tp, slide):
        """2개의 누적 막대그래프 설정"""
        chart_data = CategoryChartData()
        chart_data.categories = df_category
        chart_data.add_series("LOGICAL_READS_BLOCK", df_val1)
        chart_data.add_series("PHYSICAL_READS_BLOCK", df_val2)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED, tp["x"], tp["y"], tp["cx"], tp["cy"], chart_data
        ).chart

        return chart

    @staticmethod
    def add_stacked_three_bar_chart(df_category, df_val1, df_val2, df_val3, tp, slide):
        """3개의 누적 막대그래프 설정"""
        chart_data = CategoryChartData()
        chart_data.categories = df_category
        chart_data.add_series("LOGICAL_READS(%)", df_val1)
        chart_data.add_series("PHYSICAL_READS(%)", df_val2)
        chart_data.add_series("PHYSICAL_WRITES(%)", df_val3)

        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, tp["x"], tp["y"], tp["cx"], tp["cy"], chart_data)

    @staticmethod
    def set_chart_style(chart):
        """
        chart style
        """
        # 차트 범례 표시 및 font size, font name 설정
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        SlideManager.manuallySetLegendPosition(chart, 1, 0, 0.2, 0.2)
        chart.legend.font.name = "나눔스퀘어 네오 Light"
        chart.legend.font.size = Pt(8)

        # x축 (category) - font size, font name
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(6)
        category_axis.tick_labels.font.name = "나눔스퀘어 네오 Light"

        # X축 눈금 지우기
        category_axis.major_tick_mark = XL_TICK_MARK.NONE

        # X축 색상 회색으로 변경
        category_axis_line_format = category_axis.format.line
        category_axis_line_format.fill.solid()
        category_axis_line_format.fill.fore_color.rgb = RGBColor(191, 191, 191)

        value_axis = chart.value_axis

        # 눈금선 색상 - 회색
        gridline_format = value_axis.major_gridlines.format.line
        gridline_format.fill.solid()
        gridline_format.fill.fore_color.rgb = RGBColor(191, 191, 191)

        # y축 색상 - 흰색
        value_axis_line_format = value_axis.format.line
        value_axis_line_format.fill.solid()
        value_axis_line_format.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # y축 - font size, font name
        tick_labels = value_axis.tick_labels
        tick_labels.number_format = '#,###,"K"'
        tick_labels.font.size = Pt(7)
        tick_labels.font.name = "나눔스퀘어 네오 Light"

    @staticmethod
    def set_memory_bar_chart_style(slide, color_list):
        """memory 막대 그래프 스타일 설정"""
        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart

                plot = chart.plots[0]
                chart.has_title = False
                chart.has_legend = True
                chart.legend.include_in_layout = False
                chart.legend.position = XL_LEGEND_POSITION.TOP
                plot.vary_by_categories = False
                # chart.legend.horz_offset = 0.3
                chart.legend.font.name = "나눔스퀘어 네오 Light"
                chart.legend.font.size = Pt(9)

                legend_text = "My Legend Text"
                chart.legend.text = legend_text

                # x축 (category) - font size, font name
                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)
                category_axis.tick_labels.font.name = "맑은 고딕(본문)"
                category_axis.tick_labels.font.color.rgb = RGBColor(124, 124, 124)  # X축 색상 회색으로

                # 세로 눈금선 색상 - 회색
                height_gridline_format = category_axis.major_gridlines.format.line
                height_gridline_format.fill.solid()
                height_gridline_format.fill.fore_color.rgb = RGBColor(217, 217, 217)

                # X축 눈금 지우기
                category_axis.major_tick_mark = XL_TICK_MARK.NONE

                # X축 색상 밝은 회색 으로 변경
                category_axis_line_format = category_axis.format.line
                category_axis_line_format.fill.solid()
                category_axis_line_format.fill.fore_color.rgb = RGBColor(217, 217, 217)

                value_axis = chart.value_axis

                # 가로 눈금선 색상 - 흰색
                width_gridline_format = value_axis.major_gridlines.format.line
                width_gridline_format.fill.solid()
                width_gridline_format.fill.fore_color.rgb = RGBColor(255, 255, 255)

                # y축 색상 - 회색
                value_axis_line_format = value_axis.format.line
                value_axis_line_format.fill.solid()
                value_axis_line_format.fill.fore_color.rgb = RGBColor(217, 217, 217)

                # y축 - font size, font name
                tick_labels = value_axis.tick_labels
                tick_labels.number_format = '#"%"'
                tick_labels.font.size = Pt(9)
                tick_labels.font.name = "맑은 고딕(본문)"
                tick_labels.font.color.rgb = RGBColor(124, 124, 124)

                chart.plots[0].has_data_labels = True
                data_labels = chart.plots[0].data_labels
                data_labels.number_format = '0.00"%"'
                data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

                data_labels.font.color.rgb = RGBColor(124, 124, 124)
                data_labels.font.size = Pt(9)
                data_labels.font.name = "맑은고딕"
                data_labels.font.bold = False

                for index, series in enumerate(plot.series):
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = color_list[index]

    @staticmethod
    def set_bar_chart_style(slide):
        """막대그래프 스타일 설정"""
        color_list = [
            RGBColor(79, 162, 139),
            RGBColor(121, 107, 123),
            RGBColor(253, 107, 123),
            RGBColor(255, 192, 0),
            RGBColor(68, 114, 196),
        ]

        y_axis_title = ["count", "Avg(day)"]
        chart_list = []

        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                chart_list.append(chart)

                # x축 (category) - font size, font name
                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(9)
                category_axis.tick_labels.font.name = "나눔스퀘어 네오 Light"

                # X축 눈금 지우기
                category_axis.major_tick_mark = XL_TICK_MARK.NONE

                # X축 색상 회색으로 변경
                category_axis_line_format = category_axis.format.line
                category_axis_line_format.fill.solid()
                category_axis_line_format.fill.fore_color.rgb = RGBColor(191, 191, 191)

                value_axis = chart.value_axis

                # 눈금선 색상 - 회색
                gridline_format = value_axis.major_gridlines.format.line
                gridline_format.fill.solid()
                gridline_format.fill.fore_color.rgb = RGBColor(191, 191, 191)

                # y축 색상 - 흰색
                value_axis_line_format = value_axis.format.line
                value_axis_line_format.fill.solid()
                value_axis_line_format.fill.fore_color.rgb = RGBColor(255, 255, 255)

                # y축 - font size, font name
                tick_labels = value_axis.tick_labels
                # tick_labels.number_format = '#,###,"K"'
                tick_labels.font.size = Pt(9)
                tick_labels.font.name = "나눔스퀘어 네오 Light"

                bar_plot = chart.plots[0]
                bar_plot.gap_width = 219
                bar_plot.overlap = -27

                for index, series in enumerate(bar_plot.series):
                    for point_indx, point in enumerate(series.points):  # 차트에서 막대그래프 색상 설정
                        fill = point.format.fill
                        fill.solid()
                        fill.fore_color.rgb = color_list[point_indx]

            for index, chart in enumerate(chart_list):
                chart.value_axis.axis_title.text_frame.text = y_axis_title[index]
                chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
                chart.value_axis.axis_title.text_frame.paragraphs[0].font.bold = False

    @staticmethod
    def set_chart_legend(chart, legend_exist, layout_exist, legend_position, legend_x_position, font_name, font_size):
        """LEGEND 설정"""
        chart.has_legend = legend_exist
        chart.legend.include_in_layout = layout_exist
        chart.legend.position = legend_position
        chart.legend.horz_offset = legend_x_position
        chart.legend.font.name = font_name
        chart.legend.font.size = font_size

    @staticmethod
    def set_x_axis_category(chart, font_size, font_name, major_tick_exist, x_axis_line_color):
        """X축 설정"""
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = font_size
        category_axis.tick_labels.font.name = font_name

        # X축 눈금 존재여부 - FALSE
        category_axis.major_tick_mark = major_tick_exist

        # X축 색상 - 회색
        category_axis_line_format = category_axis.format.line
        category_axis_line_format.fill.solid()
        category_axis_line_format.fill.fore_color.rgb = x_axis_line_color

    @staticmethod
    def set_y_axis_value(
        chart, y_axis_line_color, y_axis_color, tick_label_format, tick_label_font_size, tick_label_font_name
    ):
        """Y축 설정"""
        value_axis = chart.value_axis

        # 가로 눈금선 색상 - 회색
        gridline_format = value_axis.major_gridlines.format.line
        gridline_format.fill.solid()
        gridline_format.fill.fore_color.rgb = y_axis_line_color

        # y축 색상 - 흰색
        value_axis_line_format = value_axis.format.line
        value_axis_line_format.fill.solid()
        value_axis_line_format.fill.fore_color.rgb = y_axis_color

        # y축 - font size, font name
        tick_labels = value_axis.tick_labels
        tick_labels.number_format = tick_label_format
        tick_labels.font.size = tick_label_font_size
        tick_labels.font.name = tick_label_font_name

    @staticmethod
    def set_stacked_bar_color(plot, color_list):
        """누적 막대 그래프 색상 설정"""
        for index, series in enumerate(plot.series):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = color_list[index]

    @staticmethod
    def set_stacked_bar_style(slide, tick_label_format):
        """누적 막대 그래프 스타일 서정"""
        color_list = [RGBColor(0, 176, 240), RGBColor(255, 192, 0), RGBColor(146, 208, 80)]

        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                plot = chart.plots[0]

                SlideManager.set_chart_legend(chart, True, False, XL_LEGEND_POSITION.TOP, 0.3, "나눔스퀘어 네오 Light", Pt(8))
                SlideManager.set_x_axis_category(chart, Pt(7), "맑은 고딕(본문)", XL_TICK_MARK.NONE, RGBColor(191, 191, 191))
                SlideManager.set_y_axis_value(
                    chart, RGBColor(191, 191, 191), RGBColor(255, 255, 255), tick_label_format, Pt(8), "맑은 고딕(본문)"
                )
                SlideManager.set_stacked_bar_color(plot, color_list)

    @staticmethod
    def set_label_style(label, label_font_color, label_font_size, label_font_name, label_font_bold, label_position):
        """label style"""
        label.data_label.font.color.rgb = label_font_color
        label.data_label.font.size = label_font_size
        label.data_label.font.name = label_font_name
        label.data_label.font.bold = label_font_bold
        label.data_label.position = label_position

    @staticmethod
    def set_chart_marker_style(series, marker_style, marker_fill_color, marker_line_color):
        """marker style"""
        marker = series.marker
        marker.style = marker_style
        marker.format.fill.solid()
        marker.format.fill.fore_color.rgb = marker_fill_color
        marker.format.line.color.rgb = marker_line_color

    @staticmethod
    def set_chart_line_style(series, line_color, line_width):
        """chart_line_style"""
        line = series.format.line
        line.color.rgb = line_color
        line.width = line_width

    @staticmethod
    def convert_unit(col_tuple):
        """숫자 길이 확인 후, K와 M 단위 설정"""
        new_unit_val = []

        for unit_val in col_tuple:
            if 4 <= len(str(unit_val)) < 7:
                new_val = f"{unit_val/1000:.1f}K"
                new_unit_val.append(new_val)

            elif len(str(unit_val)) >= 7:
                new_val = f"{unit_val / 1000000:.1f}M"
                new_unit_val.append(new_val)

            else:
                new_unit_val.append(str(unit_val))

        return tuple(new_unit_val)

    @staticmethod
    def set_y_axis_max_value(chart, max_value):
        """y축 value 설정"""
        chart.has_title = False

        chart_value_axis = chart.value_axis
        chart_value_axis.minimum_scale = 0
        chart_value_axis.maximum_scale = int(max_value)
        tick_labels = chart_value_axis.tick_labels

        if chart_value_axis.maximum_scale < 1000:
            chart_interval = (chart_value_axis.maximum_scale) / 5  # chart_value_axis.maximum_scale이 20일 경우, interval은 4
            chart_value_axis.major_unit = chart_interval
            tick_labels.number_format = "General"

        elif 1000 <= chart_value_axis.maximum_scale < 10000:
            chart_interval = (chart_value_axis.maximum_scale) / 5  # chart_value_axis.maximum_scale이 20일 경우, interval은 4
            chart_value_axis.major_unit = chart_interval
            tick_labels.number_format = '0.###,"K"'

        elif 10000 <= chart_value_axis.maximum_scale < 1000000:
            chart_interval = (chart_value_axis.maximum_scale) / 5  # chart_value_axis.maximum_scale이 20일 경우, interval은 4
            chart_value_axis.major_unit = chart_interval
            tick_labels.number_format = '#,###,"K"'

        # elif 1000000 <= chart_value_axis.maximum_scale < 10000000:
        #     tick_labels.number_format = '0.######,"M"'

        elif 1000000 <= chart_value_axis.maximum_scale < 100000000:
            chart_interval = (chart_value_axis.maximum_scale) / 5  # chart_value_axis.maximum_scale이 20일 경우, interval은 4
            chart_value_axis.major_unit = chart_interval
            # tick_labels.number_format = '#,######,"M"'
            tick_labels.number_format = '0,,"M"'

        return tick_labels.number_format

    @staticmethod
    def set_table_row(shape, indx, color, column_font_size):
        """table에서 특정 row cell 색상 설정 및 font_size 설정"""
        for cell in shape.table.rows[indx].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = color

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = column_font_size

    @staticmethod
    def table_cell_color(cell, color):
        """cell color 설정"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    @staticmethod
    def merge_table_cell(shape, x_1, y_1, x_2, y_2):
        """cell merge 기능"""
        shape.table.cell(x_1, y_1).merge(shape.table.cell(x_2, y_2))

    @staticmethod
    def font_alignment(specific_cell):
        """특정 cell font 가운데 정렬"""
        for paragraph in specific_cell.text_frame.paragraphs:
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    @staticmethod
    def font_center_alignment(slide):
        """table 모든 cell font 가운데 정렬"""
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    @staticmethod
    def set_font_size(shape):
        """font_size 설정"""
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(9)
                        run.font.name = "나눔스퀘어 네오 Light"

    @staticmethod
    def add_table(slide, row_count, col_count, tp, column_font_size, total_font_size):
        """table 추가 기능"""
        shape = slide.shapes.add_table(row_count, col_count, tp["left"], tp["top"], tp["width"], tp["height"])

        for row in shape.table.rows:
            for cell in row.cells:
                # 모든 셀 하얀색으로
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = total_font_size

            # # 0열 아이보리색
            cell = row.cells[0]
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

        SlideManager.set_table_row(shape, 0, RGBColor(63, 82, 108), column_font_size)  # 0행 파란색으로

    @staticmethod
    def make_table(df, tp, slide, column_font_size, total_font_size, column_width_inches=None, row_height_inches=None):
        """make table"""
        shape = df_to_table(slide, df, tp["left"], tp["top"], tp["width"], tp["height"])
        for row in shape.table.rows:
            for cell in row.cells:
                # 모든 셀 하얀색으로
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = total_font_size

            # # 0열 아이보리색
            cell = row.cells[0]
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

        SlideManager.set_table_row(shape, 0, RGBColor(63, 82, 108), column_font_size)  # 0행 파란색으로

        # value가 없는 cell은 자동으로 18포인트 크기로 형성되어 cell 크기가 망가짐 -> value가 없는 cell도 8포인트 크기로 고정
        if not cell.text_frame.text:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = total_font_size

        if column_width_inches is not None:
            for idx, column in enumerate(shape.table.columns):
                column.width = column_width_inches[idx]

        if row_height_inches is not None:
            for idx, row in enumerate(shape.table.rows):
                row.height = row_height_inches[idx]

    @staticmethod
    def adjust_row_column_size(slide, column_width_inches=None, row_height_inches=None):
        """row_column_size 조정"""
        for shape in slide.shapes:
            if shape.has_table:
                if column_width_inches is not None:
                    for idx, column in enumerate(shape.table.columns):
                        column.width = column_width_inches[idx]

                if row_height_inches is not None:
                    for idx, row in enumerate(shape.table.rows):
                        row.height = row_height_inches[idx]

    @staticmethod
    def activity_info_style(slide, val):
        """activity_info style 설정"""
        text1 = (
            "성능 분석 대상 및 구간 선정\n"
            "성능 관련 항목 점검(MaxGauge, Dynamic Performance View, Automatic Workload Repository)\n"
            "성능 안전점검 및 성능개선을 위한 분석 진행\n"
            "분석 결과에 따른 권고사항 제시"
        )

        text2 = "성능관련 이슈에 대한 사전예방\n" "개선안 적용으로 보다 안정적이고 효율적인 시스템 운영"

        for shape in slide.shapes:
            if shape.has_table:
                shape.table.cell(0, 0).text = "구분"
                shape.table.cell(0, 1).text = "설명"
                shape.table.cell(1, 0).text = "목표"
                shape.table.cell(1, 1).text = "DBMS 및 OS, 시스템, 어플리케이션의 이해를 바탕으로, 성능저하 요소를 찾아 제거/수정하여 성능을 개선"
                shape.table.cell(2, 0).text = "대상"
                shape.table.cell(2, 1).text = val
                shape.table.cell(3, 0).text = "주요 활동"
                shape.table.cell(3, 1).text = text1
                shape.table.cell(4, 0).text = "기대 효과"
                shape.table.cell(4, 1).text = text2

                # font size 9로 설정
                SlideManager.set_font_size(shape)

                # 특정 셀 font size 11로 설정
                shape.table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                shape.table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.size = Pt(11)

                # 특정 셀 가운데 정렬
                SlideManager.font_alignment(shape.table.cell(0, 1))
                for i in range(0, 5):
                    SlideManager.font_alignment(shape.table.cell(i, 0))

    @staticmethod
    def schedule_style(slide):
        """schedule style 설정"""
        text1 = "MaxGauge\n" "Dynamic Performance View\n" "Automatic Workload Repository\n" "기타 Script\n"

        for shape in slide.shapes:
            if shape.has_table:
                shape.table.cell(0, 0).text = "구분"
                shape.table.cell(0, 1).text = "설명"
                shape.table.cell(1, 0).text = "수행인력"
                shape.table.cell(1, 1).text = "엑셈 DB 기술 본부"

                shape.table.cell(2, 1).text = "소속 및 직급"
                shape.table.cell(2, 2).text = "성명"
                shape.table.cell(2, 3).text = "역할"

                shape.table.cell(3, 1).text = "-"
                shape.table.cell(3, 2).text = "-"
                shape.table.cell(3, 3).text = "분석, 개선 및 보고"

                shape.table.cell(4, 0).text = "수행일정"
                shape.table.cell(5, 0).text = "데이터 수집 및 분석"
                shape.table.cell(5, 1).text = text1
                SlideManager.table_cell_color(shape.table.cell(1, 1), RGBColor(242, 242, 242))

        for shape in slide.shapes:
            if shape.has_table:
                SlideManager.merge_table_cell(shape, 0, 1, 0, 3)
                SlideManager.merge_table_cell(shape, 1, 1, 1, 3)
                SlideManager.merge_table_cell(shape, 4, 1, 4, 3)
                SlideManager.merge_table_cell(shape, 5, 1, 5, 3)
                SlideManager.merge_table_cell(shape, 1, 0, 3, 0)

            SlideManager.set_font_size(shape)

            # 특정 셀 font size 11로 설정
            shape.table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            shape.table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.size = Pt(11)

            # 가운데 정렬
            SlideManager.font_alignment(shape.table.cell(0, 0))
            SlideManager.font_alignment(shape.table.cell(0, 1))
            SlideManager.font_alignment(shape.table.cell(1, 0))
            SlideManager.font_alignment(shape.table.cell(1, 1))
            SlideManager.font_alignment(shape.table.cell(2, 1))
            SlideManager.font_alignment(shape.table.cell(2, 2))
            SlideManager.font_alignment(shape.table.cell(4, 0))
            SlideManager.font_alignment(shape.table.cell(5, 0))

    @staticmethod
    def time_model_table_style(slide, color):
        """time_model_table style 설정"""
        for shape in slide.shapes:
            if shape.has_table:
                for i in range(13, 17):
                    SlideManager.set_table_row(shape, i, color, Pt(8))

    @staticmethod
    def db_system_table_style(slide):
        """db_system_table style 설정"""
        for shape in slide.shapes:
            if shape.has_table:
                shape.table.cell(8, 0).text = "SGA"
                shape.table.cell(8, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(
                    8
                )  # Set the desired font size (e.g., Pt(12))
                SlideManager.merge_table_cell(shape, 8, 0, 12, 0)
                SlideManager.merge_table_cell(shape, 13, 0, 13, 1)
                for i in range(0, 8):
                    SlideManager.merge_table_cell(shape, i, 0, i, 1)

                for i in range(8, 13):
                    cell_position = shape.table.cell(i, 1)
                    SlideManager.table_cell_color(cell_position, RGBColor(218, 227, 243))

                for row in shape.table.rows:
                    for cell in row.cells:
                        SlideManager.font_alignment(cell)

    @staticmethod
    def load_profile_style(slide):
        """load_profile style 설정"""
        for shape in slide.shapes:
            if shape.has_table:
                for i in range(1, 16, 2):
                    SlideManager.merge_table_cell(shape, i, 0, i + 1, 0)
                    SlideManager.merge_table_cell(shape, i, 1, i + 1, 1)

    @staticmethod
    def optimizer_statistics_style(slide):
        """optimiser_statistics style 설정"""
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row_number, row in enumerate(shape.table.rows):
                    if row.cells[5].text == "SUCCEEDED":
                        run = row.cells[5].text_frame.paragraphs[0].runs[0]
                        run.font.color.rgb = RGBColor(0, 112, 192)
                        run.font.bold = True
                        # row.cells[5].text_frame.paragraphs[0].runs[0].font.color = RGBColor(0,112,192)

                    elif row.cells[5].text == "FAILED":
                        run = row.cells[5].text_frame.paragraphs[0].runs[0]
                        run.font.color.rgb = RGBColor(255, 0, 0)
                        run.font.bold = True
                        # row.cells[5].text_frame.paragraphs[0].runs[0].font.color = RGBColor(255,0,0)

                    if row_number > 0 and row_number % 2 == 0:
                        for cell in row.cells:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                    if row_number > 0 and row_number % 2 != 0:
                        for cell in row.cells:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    @staticmethod
    def db_parameter_style(df, tp, slide, column_width_inches):
        """db_parameter style 설정"""
        shape = df_to_table(slide, df, tp["left"], tp["top"], tp["width"], tp["height"])
        for idx, row in enumerate(shape.table.rows):
            row.height = Inches(0.287)

            if idx > 0:
                run = row.cells[4].text_frame.paragraphs[0].runs[0]
                run.font.color.rgb = RGBColor(255, 59, 59)
                run.font.bold = True

            for cell in row.cells:
                # 모든 셀 하얀색으로
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(9)

            # # 0열 아이보리색
            cell = row.cells[0]
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(157, 203, 190)

        SlideManager.set_table_row(shape, 0, RGBColor(99, 173, 152))  # 0행 파란색으로

        if column_width_inches is not None:
            for idx, column in enumerate(shape.table.columns):
                column.width = column_width_inches[idx]

        # value가 없는 cell은 자동으로 18포인트 크기로 형성되어 cell 크기가 망가짐 -> value가 없는 cell도 8포인트 크기로 고정
        if not cell.text_frame.text:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)
                    run.font.size = Pt(8)

    @staticmethod
    def db_parameter_cell_merge(slide):
        """db_parameter_table cell merge 설정"""
        for shape in slide.shapes:
            if shape.has_table:
                row_1 = shape.table.rows[1]
                row_last = shape.table.rows[len(shape.table.rows) - 1]

                cell_1 = row_1.cells[0]
                cell_2 = row_last.cells[0]
                cell_1.merge(cell_2)

    @staticmethod
    def SubElement(parent, tagname, **kwargs):
        """db_parameter style 설정"""
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    @staticmethod
    def manuallySetLegendPosition(chart, x, y, w, h):
        """db_parameter style 설정"""
        L = chart.legend._element.get_or_add_layout()
        mL = L.get_or_add_manualLayout()

        SlideManager.SubElement(mL, "c:xMode", val="edge")
        SlideManager.SubElement(mL, "c:yMode", val="edge")

        SlideManager.SubElement(mL, "c:x", val=str(x))
        SlideManager.SubElement(mL, "c:y", val=str(y))
        SlideManager.SubElement(mL, "c:w", val=str(w))
        SlideManager.SubElement(mL, "c:h", val=str(h))

    @staticmethod
    def database_style(slide, df):
        """database parameter style 설정"""
        same_row = df[df["권장값"] == df["현재값"]]  # 권장값과 현재값이 같은 row

        tuple_list = []
        result_dict = {}
        for shape in slide.shapes:
            if shape.has_table:
                for row_number, row in enumerate(shape.table.rows):
                    row.height = Inches(0.4012)
                    if row_number != 0:
                        if row_number in same_row.index + 1:
                            for cell in row.cells:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(
                                    217, 217, 217
                                )  # 권장값 false, 현재값 false에 해당하는 cell 밝은 회색으로

                        else:
                            target_cell = shape.table.cell(row_number, 4)
                            run = target_cell.text_frame.paragraphs[0].runs[0]
                            run.font.color.rgb = RGBColor(255, 59, 59)  # 권장값 글씨 빨간색으로
                            run.font.bold = True

                        zero_col_cell = row.cells[0]
                        zero_col_cell.fill.solid()
                        zero_col_cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

                        for paragraph in zero_col_cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)
                            paragraph.font.name = "나눔스퀘어"
                            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # 0행을 제외한 0열에 해당하는 값들 가운데 정렬

                        if zero_col_cell.text != "":
                            start_num = row_number

                        # if start_num < row_number and zero_col_cell.text == "":
                        if zero_col_cell.text == "":
                            tuple_list.append((start_num, row_number))

                for x, y in tuple_list:
                    if x in result_dict:
                        result_dict[x] = max(result_dict[x], y)
                    else:
                        result_dict[x] = y

                result_list = [(key, value) for key, value in result_dict.items()]

                for tuple_val in result_list:
                    a, b = tuple_val
                    SlideManager.merge_table_cell(shape, a, 0, b, 0)

                # 가운데 정렬
                for i in range(0, 6):
                    SlideManager.font_alignment(shape.table.cell(0, i))
